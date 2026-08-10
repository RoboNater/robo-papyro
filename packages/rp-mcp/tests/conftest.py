"""Fixtures for the rp-mcp suite: a sandboxed directory tree, sample documents,
and the in-process MCP client the server tests drive.

**No subprocess, no sockets.** ``mcp.Client`` accepts an ``MCPServer`` object
and connects to it over in-memory streams, so a tool test exercises the real
registration, the real JSON-schema validation, and the real error path without
a transport in the way. That is also why these tests are worth having on top of
the leaves' own: they check the *wiring*, and the wiring is all this package is.

Documents are generated here, exactly as in the leaves' conftests — no binary
fixtures in git.
"""

from __future__ import annotations

import json
import shutil
from pathlib import Path
from typing import Any

import anyio
import pytest
from mcp import Client
from mcp.server.mcpserver import MCPServer
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas as rl_canvas

import rp_docx
import rp_pptx
import rp_xlsx
from rp_mcp.sandbox import ROOTS_ENV, WRITE_ROOT_ENV, Sandbox

PDF_PAGES = ["Alpha page one", "Beta page two", "Gamma page three"]

DOCX_MARKDOWN = """\
# Quarterly Report

An opening paragraph about the quarter.

## Findings

- Revenue rose
- Costs fell

| Region | Total |
|--------|-------|
| North  | 12    |
| South  | 8     |
"""

PPTX_MARKDOWN = """\
# Quarterly Deck

## Findings

- Revenue rose
- Costs fell

## Next Steps

- Hire
- Ship
"""

TEMPLATE_MARKDOWN = """\
# Engagement Letter

Dear {{ client_name }},

We are pleased to confirm the engagement dated {{ start_date }}.
"""


def poppler_available() -> bool:
    return bool(shutil.which("pdftotext") and shutil.which("pdftoppm"))


def pytest_collection_modifyitems(config, items):
    """Skip @pytest.mark.requires_poppler tests when poppler is absent.

    Same shape as rp-pdf's, and for the same reason: a marker rather than an
    importable skipif, because pytest's importlib import mode keeps a tests
    directory off ``sys.path``.
    """
    if poppler_available():
        return
    skip = pytest.mark.skip(reason="poppler (pdftoppm/pdftotext) not installed")
    for item in items:
        if "requires_poppler" in item.keywords:
            item.add_marker(skip)


@pytest.fixture(autouse=True)
def clean_env(monkeypatch: pytest.MonkeyPatch) -> None:
    """Neither sandbox variable leaks in from the machine running the tests.

    ``RP_MCP_ROOTS`` set in a developer's shell would silently widen every
    sandbox under test, which is the one thing these tests exist to pin down.
    """
    monkeypatch.delenv(ROOTS_ENV, raising=False)
    monkeypatch.delenv(WRITE_ROOT_ENV, raising=False)


@pytest.fixture
def docs(tmp_path: Path) -> Path:
    """The directory a server is pointed at."""
    path = tmp_path / "docs"
    path.mkdir()
    return path


@pytest.fixture
def outbox(tmp_path: Path) -> Path:
    """The directory a writable server may write into."""
    path = tmp_path / "outbox"
    path.mkdir()
    return path


@pytest.fixture
def outside(tmp_path: Path) -> Path:
    """A file the server was never pointed at, for the containment tests."""
    path = tmp_path / "outside" / "secret.docx"
    path.parent.mkdir()
    path.write_text("not for the server", encoding="utf-8")
    return path


@pytest.fixture
def sample_pdf(docs: Path) -> Path:
    path = docs / "report.pdf"
    canvas = rl_canvas.Canvas(str(path), pagesize=letter)
    for text in PDF_PAGES:
        canvas.drawString(72, 720, text)
        canvas.showPage()
    canvas.save()
    return path


@pytest.fixture
def sample_docx(docs: Path) -> Path:
    path = docs / "report.docx"
    rp_docx.create(path, markdown=DOCX_MARKDOWN, title="Quarterly Report")
    return path


@pytest.fixture
def sample_pptx(docs: Path) -> Path:
    path = docs / "deck.pptx"
    rp_pptx.create(path, markdown=PPTX_MARKDOWN)
    return path


@pytest.fixture
def pdf_with_images(docs: Path, tmp_path: Path) -> Path:
    """A four-page PDF carrying one embedded image per page.

    Exists for the accumulate-into-one-directory tests: rp-pdf names extracted
    files per page, and that is what makes `--pages 1-2` then `3-4` into a
    single folder safe. Proving it needs a document with images on known pages.
    """
    photo = tmp_path / "photo.png"
    PILImage.new("RGB", (32, 24), "red").save(photo)
    path = docs / "illustrated.pdf"
    canvas = rl_canvas.Canvas(str(path), pagesize=letter)
    for number in range(1, 5):
        canvas.drawImage(ImageReader(str(photo)), 72, 600, width=32, height=24)
        canvas.drawString(72, 720, f"page {number}")
        canvas.showPage()
    canvas.save()
    return path


@pytest.fixture
def pptx_with_images(docs: Path, tmp_path: Path) -> Path:
    """A four-slide deck with one picture per slide.

    The rp-pptx counterpart of :func:`pdf_with_images`. rp-pptx numbers images
    across the whole deck *before* applying the slide filter, so a second range
    continues the numbering instead of restarting it — which is the claim these
    tests exist to hold up.
    """
    photo = tmp_path / "photo.png"
    PILImage.new("RGB", (32, 24), "blue").save(photo)
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for _ in range(4):
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(str(photo), Inches(1), Inches(1))
    path = docs / "illustrated.pptx"
    presentation.save(path)
    return path


@pytest.fixture
def docx_template(docs: Path) -> Path:
    """A document with two `{{ placeholder }}` names in it."""
    path = docs / "engagement.docx"
    rp_docx.create(path, markdown=TEMPLATE_MARKDOWN)
    return path


@pytest.fixture
def pptx_template(docs: Path) -> Path:
    path = docs / "engagement.pptx"
    rp_pptx.create(path, markdown="# {{ client_name }}\n\n## Terms\n\n- Dated {{ start_date }}\n")
    return path


@pytest.fixture
def sample_xlsx(docs: Path) -> Path:
    from rp_xlsx.models import SheetSpec

    path = docs / "report.xlsx"
    rp_xlsx.create(
        path,
        sheets=[
            SheetSpec(
                name="Data",
                header=["Region", "Total"],
                rows=[["North", 12], ["South", 8]],
            )
        ],
    )
    return path


@pytest.fixture
def xlsx_with_images(docs: Path, tmp_path: Path) -> Path:
    """A four-sheet workbook with one picture per sheet.

    The rp-xlsx counterpart of :func:`pptx_with_images`: images are numbered
    across the whole workbook before the sheet filter is applied, so a second
    range continues the numbering instead of restarting it.
    """
    import openpyxl
    from openpyxl.drawing.image import Image as XlImage

    photo = tmp_path / "photo.png"
    PILImage.new("RGB", (32, 24), "blue").save(photo)
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for number in range(1, 5):
        ws = wb.create_sheet(f"Sheet{number}")
        ws["A1"] = f"sheet {number}"
        ws.add_image(XlImage(str(photo)), "C1")
    path = docs / "illustrated.xlsx"
    wb.save(path)
    return path


@pytest.fixture
def xlsx_template(docs: Path) -> Path:
    import openpyxl

    path = docs / "engagement.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Letter"
    ws["A1"] = "Dear {{ client_name }},"
    ws["A2"] = "We are pleased to confirm the engagement dated {{ start_date }}."
    wb.save(path)
    return path


@pytest.fixture
def read_sandbox(docs: Path) -> Sandbox:
    """A read-only server's sandbox: one root, no write root."""
    return Sandbox([docs])


@pytest.fixture
def write_sandbox(docs: Path, outbox: Path) -> Sandbox:
    """A writable server's sandbox."""
    return Sandbox([docs], write_root=outbox)


# --- driving a server in-process -------------------------------------------


class Driver:
    """Talks to an ``MCPServer`` the way a client does, synchronously.

    Delivered as the ``mcp`` fixture rather than imported: pytest runs in
    importlib import mode, so a test module's directory is not on ``sys.path``
    and ``from conftest import ...`` does not work (see AGENTS.md).
    """

    @staticmethod
    def call(server: MCPServer, name: str, arguments: dict[str, Any] | None = None):
        """One tool call over in-memory streams.

        A fresh session per call: tool handlers hold no state between calls, and
        an isolated session keeps a failure in one test from becoming a mystery
        in the next.
        """

        async def run():
            async with Client(server) as client:
                return await client.call_tool(name, arguments or {})

        return anyio.run(run)

    @staticmethod
    def listed(server: MCPServer) -> list[Any]:
        """The tools a client sees in ``tools/list``, sorted by name."""

        async def run():
            async with Client(server) as client:
                return (await client.list_tools()).tools

        return sorted(anyio.run(run), key=lambda tool: tool.name)

    def names(self, server: MCPServer) -> set[str]:
        return {tool.name for tool in self.listed(server)}

    def schema(self, server: MCPServer, name: str) -> dict:
        return next(tool.input_schema for tool in self.listed(server) if tool.name == name)

    @staticmethod
    def text(result) -> str:
        return "\n".join(block.text for block in result.content if getattr(block, "text", None))

    def structured(self, result) -> Any:
        """A successful call's structured content.

        Fails loudly on an error result rather than returning ``None``: a test
        that silently compares ``None`` to ``None`` is exactly the shape
        AGENTS.md warns about.
        """
        assert not result.is_error, self.text(result)
        assert result.structured_content is not None, "tool returned no structured content"
        return result.structured_content

    def envelope(self, result) -> dict:
        """The error envelope from a failed call: the *last* line of its text.

        The ordering is the contract — human message first, envelope last — so
        this reads the last line rather than hunting for a brace.
        """
        assert result.is_error, "expected the call to fail"
        return json.loads(self.text(result).splitlines()[-1])

    def error_type(self, result) -> str:
        return self.envelope(result)["error"]["type"]


@pytest.fixture
def mcp() -> Driver:
    return Driver()


@pytest.fixture
def elsewhere(tmp_path: Path) -> str:
    """An absolute path under no root at all."""
    return str(tmp_path / "elsewhere" / "file.bin")
