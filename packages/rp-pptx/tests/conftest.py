"""Fixtures: every deck and template used by the suite, built at run time.

**Nothing binary in git** (spec section 11.1). A corporate `.potx` in the
repository is a licensing question, an opaque diff, and a debugging hazard at
once — when a test fails you cannot tell whether the code or the template
changed. Everything here is constructed, so a failure is always the code.

The three templates are **adversarial rather than realistic**, because realism is
not what catches bugs here (spec section 11.2):

``minimal``
    python-pptx's bundled default: stock layout names, 4:3. The happy path, the
    default-``LayoutMap`` path, and both sides of the aspect contract.
``house_like``
    What a real house template looks like: renamed layouts, a name carrying a
    space and a non-ASCII character, 16:9, an image on the master, a second
    master, and a ``.layoutmap.json`` beside it.
``hostile``
    Missing two layouts the default ``LayoutMap`` maps to, no layoutmap, and two
    layout names differing only by case. Exists to prove failures are loud, and
    to make lazy layout checking observable: the roles it *does* have still work.

**On the modern-comments fixture.** Spec section 7 defers modern threaded
comments until a real PowerPoint-authored deck can be inspected, and section 11.1
is explicit that the generator must be written from that file rather than from
what the schema implies. No such file was available, so
:func:`modern_comments_deck` does not claim to reproduce PowerPoint's markup. It
writes a part with the modern *content type* and nothing else that is trustworthy
— which is exactly and only what the deferral path keys on. The tests built on it
assert the deferral behaves (an ``UnsupportedFeatureError`` naming the slides,
``comment_count`` going null), never that the comment bodies parse.
"""

from __future__ import annotations

import functools
import io
import json
import shutil
import subprocess
import tempfile
from pathlib import Path

import pytest
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu, Inches, Pt

from rp_core import ooxml as core_ooxml
from rp_pptx import ooxml

#: Text that must never appear in a manifest built from ``house_like`` — the
#: redaction assertion in spec section 5.2 searches for it.
SECRET_TEXT = "Zephyr Quokka Confidential Boilerplate"

#: The house layout names. Deliberately none of PowerPoint's own: a real house
#: template renames everything, which is why LayoutMap exists.
HOUSE_LAYOUTS = {
    "title": "RP Title",
    "section": "House Section Break",
    "content": "House Content",
    "blank": "House Blank",
}

#: A layout name with a space and a non-ASCII character (spec section 11.2).
ACCENTED_LAYOUT = "Résumé Layout"

IMAGE_SIZE = (48, 32)


# --- LibreOffice, probed rather than assumed --------------------------------


@functools.lru_cache(maxsize=1)
def soffice_available() -> bool:
    """Whether LibreOffice is present **and can actually convert a deck**.

    Checking only for the binary is not enough — a container can ship `soffice`
    that fails every conversion for want of a filter or a Java runtime. A
    presence check turns that into confusing failures in an environment that
    never claimed to support conversion; a functional probe turns it into honest
    skips. The pattern is rp-docx's, and spec section 11.3 asks for it by name.
    """
    if shutil.which("soffice") is None:
        return False
    with tempfile.TemporaryDirectory(prefix="rp-pptx-probe-") as tmp:
        directory = Path(tmp)
        source = directory / "probe.pptx"
        Presentation().save(str(source))
        try:
            subprocess.run(
                [
                    "soffice",
                    f"-env:UserInstallation={(directory / 'profile').as_uri()}",
                    "--headless",
                    "--norestore",
                    "--invisible",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(directory),
                    str(source),
                ],
                capture_output=True,
                timeout=180,
            )
        except (OSError, subprocess.TimeoutExpired):
            return False
        return (directory / "probe.pdf").is_file()


def pytest_collection_modifyitems(config, items):
    """Skip ``@pytest.mark.requires_soffice`` tests when LibreOffice cannot work.

    A marker rather than an importable ``skipif``, because ``from conftest
    import ...`` only works when the tests directory happens to be on sys.path,
    which pytest's importlib import mode deliberately stops doing.
    """
    if not any("requires_soffice" in item.keywords for item in items):
        return
    if soffice_available():
        return
    skip = pytest.mark.skip(reason="LibreOffice (soffice) is absent or cannot convert decks")
    for item in items:
        if "requires_soffice" in item.keywords:
            item.add_marker(skip)


# --- small helpers -----------------------------------------------------------


def png_bytes(size: tuple[int, int] = IMAGE_SIZE, colour: tuple[int, int, int] = (200, 30, 30)):
    """An in-memory PNG. Returned as a stream because that is what python-pptx
    and the master-picture helper both take."""
    buffer = io.BytesIO()
    Image.new("RGB", size, colour).save(buffer, "PNG")
    buffer.seek(0)
    return buffer


def _add_master_textbox(presentation: Presentation, text: str) -> None:
    """Put a plain text box on the first master.

    ``master.shapes`` has no ``add_textbox`` — only slide shape trees do — so
    this is hand-built, the same way :func:`rp_pptx.ooxml.add_master_picture` is.
    It exists to give the redaction test something to look for: boilerplate on
    the master is exactly the kind of confidential text a real template carries.
    """
    tree = presentation.slide_masters[0].shapes._spTree
    shape = etree.SubElement(tree, ooxml.qn("p:sp"))
    nv = etree.SubElement(shape, ooxml.qn("p:nvSpPr"))
    name = etree.SubElement(nv, ooxml.qn("p:cNvPr"))
    name.set("id", "1024")
    name.set("name", "Boilerplate")
    etree.SubElement(nv, ooxml.qn("p:cNvSpPr"))
    etree.SubElement(nv, ooxml.qn("p:nvPr"))
    etree.SubElement(shape, ooxml.qn("p:spPr"))
    body = etree.SubElement(shape, ooxml.qn("p:txBody"))
    etree.SubElement(body, ooxml.qn("a:bodyPr"))
    etree.SubElement(body, ooxml.qn("a:lstStyle"))
    paragraph = etree.SubElement(body, ooxml.qn("a:p"))
    run = etree.SubElement(paragraph, ooxml.qn("a:r"))
    etree.SubElement(run, ooxml.qn("a:t")).text = text


def _bulleted(slide, items: list[tuple[int, str]]) -> None:
    """Fill a slide's body placeholder with ``(level, text)`` bullets."""
    body = next(p for p in slide.placeholders if p.placeholder_format.idx != 0)
    frame = body.text_frame
    frame.clear()
    for offset, (level, text) in enumerate(items):
        paragraph = frame.paragraphs[0] if offset == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level


# --- the three templates -----------------------------------------------------


@pytest.fixture(scope="session")
def bundled_default() -> Path:
    """python-pptx's own bundled template, by path.

    Spec section 4 turns on being able to pass *this* file explicitly and get
    different behaviour from passing ``None``, so the test that proves it needs
    the real path rather than a copy.
    """
    import pptx

    return Path(pptx.__file__).parent / "templates" / "default.pptx"


@pytest.fixture(scope="session")
def templates_dir(tmp_path_factory) -> Path:
    """A directory laid out the way ``RP_PPTX_TEMPLATE_DIR`` expects."""
    return tmp_path_factory.mktemp("templates")


@pytest.fixture(scope="session")
def minimal_template(templates_dir: Path, bundled_default: Path) -> Path:
    """The bundled default, copied under a name resolution can find. 4:3."""
    target = templates_dir / "minimal.potx"
    ooxml.retype_as_template(bundled_default, target)
    return target


@pytest.fixture(scope="session")
def house_like_template(templates_dir: Path) -> Path:
    """A stand-in for a real house template (spec section 11.2).

    Renamed layouts, one of them accented, 16:9, a logo on the master, a second
    master, confidential-looking boilerplate, and a layoutmap beside it.
    """
    target = templates_dir / "house_like.potx"
    with tempfile.TemporaryDirectory(prefix="rp-pptx-house-") as tmp:
        staged = Path(tmp) / "staged.pptx"
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        ooxml.add_master_picture(presentation, png_bytes(), name="House Logo")
        _add_master_textbox(presentation, SECRET_TEXT)
        presentation.core_properties.author = SECRET_TEXT
        presentation.core_properties.title = SECRET_TEXT
        presentation.save(str(staged))

        rebuilt = Path(tmp) / "rebuilt.pptx"
        ooxml.rebuild_masters(
            staged,
            rebuilt,
            [
                [
                    (
                        HOUSE_LAYOUTS["title"],
                        [(0, "ctrTitle", "Title 1"), (1, "subTitle", "Sub 2")],
                    ),
                    (HOUSE_LAYOUTS["content"], [(0, "title", "Title 1"), (1, "body", "Content 2")]),
                    (HOUSE_LAYOUTS["section"], [(0, "title", "Title 1")]),
                    (HOUSE_LAYOUTS["blank"], []),
                    (ACCENTED_LAYOUT, [(0, "title", "Title 1"), (1, "body", "Content 2")]),
                ],
                [("Appendix Content", [(0, "title", "Title 1"), (1, "body", "Content 2")])],
            ],
        )
        ooxml.retype_as_template(rebuilt, target)

    (templates_dir / "house_like.potx.layoutmap.json").write_text(
        json.dumps(HOUSE_LAYOUTS, indent=2), encoding="utf-8"
    )
    return target


@pytest.fixture(scope="session")
def hostile_template(templates_dir: Path) -> Path:
    """Missing layouts, no layoutmap, and a case collision (spec section 11.2).

    It keeps the *title* and *blank* roles the default ``LayoutMap`` names and
    drops *content* and *section*. That asymmetry is the point: lazy checking
    (section 5.1) means a deck needing only a title slide must still succeed
    here, and only a deck that reaches for content is allowed to fail.
    """
    target = templates_dir / "hostile.potx"
    with tempfile.TemporaryDirectory(prefix="rp-pptx-hostile-") as tmp:
        staged = Path(tmp) / "staged.pptx"
        Presentation().save(str(staged))
        rebuilt = Path(tmp) / "rebuilt.pptx"
        ooxml.rebuild_masters(
            staged,
            rebuilt,
            [
                [
                    ("Title Slide", [(0, "ctrTitle", "Title 1"), (1, "subTitle", "Sub 2")]),
                    ("Blank", []),
                    # Two names differing only by case, so anything matching
                    # case-insensitively picks one of them arbitrarily.
                    ("Overview", [(0, "title", "Title 1"), (1, "body", "Content 2")]),
                    ("overview", [(0, "title", "Title 1"), (1, "body", "Content 2")]),
                ]
            ],
        )
        ooxml.retype_as_template(rebuilt, target)
    return target


@pytest.fixture(scope="session")
def all_templates(minimal_template, house_like_template, hostile_template) -> Path:
    """All three present in one directory, for resolution tests."""
    return minimal_template.parent


@pytest.fixture
def template_env(all_templates: Path, monkeypatch) -> Path:
    """``RP_PPTX_TEMPLATE_DIR`` pointed at the template directory."""
    monkeypatch.setenv("RP_PPTX_TEMPLATE_DIR", str(all_templates))
    monkeypatch.delenv("RP_PPTX_TEMPLATE", raising=False)
    return all_templates


# --- decks -------------------------------------------------------------------


@pytest.fixture
def simple_deck(tmp_path: Path) -> Path:
    """Three titled slides with flat bullets. The workhorse for slide ops."""
    presentation = Presentation()
    for title, body in [
        ("Alpha", "first body"),
        ("Beta", "second body"),
        ("Gamma", "third body"),
    ]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        _bulleted(slide, [(0, body)])
    target = tmp_path / "simple.pptx"
    presentation.save(str(target))
    return target


@pytest.fixture
def rich_deck(tmp_path: Path) -> Path:
    """One deck carrying every shape the readers have to cope with.

    Slide 1  title + nested bullets + speaker notes
    Slide 2  a table with a merged origin cell
    Slide 3  a picture with real alt text, and a grouped shape holding text
    Slide 4  a bar chart
    """
    presentation = Presentation()

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Outline"
    _bulleted(
        slide, [(0, "top level"), (1, "nested once"), (2, "nested twice"), (0, "back to top")]
    )
    slide.notes_slide.notes_text_frame.text = "Speaker notes for the outline slide"

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Numbers"
    table = slide.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(6), Inches(2)).table
    # Merge *before* filling. python-pptx's merge concatenates the swallowed
    # cell's text onto the origin, so filling first would build a fixture whose
    # origin reads "origin\nspanned" — an artefact of how it was made rather
    # than what PowerPoint writes, and the assertion would then encode the
    # artefact.
    table.cell(0, 0).merge(table.cell(0, 1))
    values = [["origin", None, "c3"], ["r2c1", "r2c2", "r2c3"], ["r3c1", "r3c2", "r3c3"]]
    for row_index, row in enumerate(values):
        for col_index, value in enumerate(row):
            if value is not None:
                table.cell(row_index, col_index).text = value

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture = slide.shapes.add_picture(png_bytes(), Inches(1), Inches(1), Inches(2), Inches(1.3))
    picture._element._nvXxPr.cNvPr.set("descr", "A red rectangle")
    group = slide.shapes.add_group_shape()
    inner = group.shapes.add_textbox(Inches(4), Inches(1), Inches(3), Inches(1))
    inner.text_frame.text = "text inside a group"

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Chart"
    data = CategoryChartData()
    data.categories = ["East", "West"]
    data.add_series("Revenue", (11.0, 22.0))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(6), Inches(4), data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Revenue by region"

    target = tmp_path / "rich.pptx"
    presentation.save(str(target))
    return target


@pytest.fixture
def runs_deck(tmp_path: Path) -> Path:
    """A deck whose text is deliberately split across run boundaries.

    python-pptx writes a paragraph as one run, so every case spec section 6
    requires has to be built by hand: the whole point is that a placeholder can
    straddle runs, and a fixture that never splits one tests nothing.
    """
    presentation = Presentation()

    def split_paragraph(frame, parts: list[tuple[str, bool]]) -> None:
        """Replace ``frame``'s first paragraph with one run per ``(text, bold)``."""
        paragraph = frame.paragraphs[0]
        for run in list(paragraph.runs):
            run._r.getparent().remove(run._r)
        for text, bold in parts:
            run = paragraph.add_run()
            run.text = text
            run.font.bold = bold

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Runs"
    body = next(p for p in slide.placeholders if p.placeholder_format.idx != 0)
    body.text_frame.clear()
    # Split across three runs, and across a formatting boundary.
    split_paragraph(
        body.text_frame, [("Hello {{ na", False), ("me }} and ", True), ("{{ role }} here", False)]
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Table and group"
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(4), Inches(1)).table
    split_paragraph(table.cell(0, 0).text_frame, [("cell {{ na", False), ("me }}", False)])
    group = slide.shapes.add_group_shape()
    inner = group.shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
    split_paragraph(inner.text_frame, [("group {{ na", False), ("me }}", False)])
    slide.notes_slide.notes_text_frame.text = "notes mention {{ name }} too"

    target = tmp_path / "runs.pptx"
    presentation.save(str(target))
    return target


@pytest.fixture
def overlap_deck(tmp_path: Path) -> Path:
    """One slide whose text contains overlapping replacement candidates."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Hello World and Hello"
    target = tmp_path / "overlap.pptx"
    presentation.save(str(target))
    return target


# --- comment fixtures, built by adding parts python-pptx cannot ---------------

_CLASSIC_COMMENTS_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.comments+xml"
_CLASSIC_AUTHORS_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml"
)


def _classic_authors() -> bytes:
    root = etree.Element(ooxml.qn("p:cmAuthorLst"), nsmap={"p": ooxml.NS["p"]})
    for author_id, (name, initials) in enumerate([("Ada Lovelace", "AL"), ("Grace Hopper", "GH")]):
        author = etree.SubElement(root, ooxml.qn("p:cmAuthor"))
        author.set("id", str(author_id))
        author.set("name", name)
        author.set("initials", initials)
        author.set("lastIdx", "1")
        author.set("clrIdx", str(author_id))
    return ooxml.serialize(root)


def _classic_comments(entries: list[tuple[int, str, str]]) -> bytes:
    """``(author_id, date, text)`` triples as one slide's classic comment part."""
    root = etree.Element(ooxml.qn("p:cmLst"), nsmap={"p": ooxml.NS["p"]})
    for index, (author_id, date, text) in enumerate(entries, start=1):
        comment = etree.SubElement(root, ooxml.qn("p:cm"))
        comment.set("authorId", str(author_id))
        comment.set("dt", date)
        comment.set("idx", str(index))
        position = etree.SubElement(comment, ooxml.qn("p:pos"))
        position.set("x", "100")
        position.set("y", "100")
        etree.SubElement(comment, ooxml.qn("p:text")).text = text
    return ooxml.serialize(root)


def _add_parts(
    source: Path,
    target: Path,
    parts: dict[str, bytes],
    overrides: dict[str, str],
    slide_rels: dict[int, list[tuple[str, str]]],
    presentation_rels: list[tuple[str, str]],
) -> Path:
    """Add parts, content-type overrides, and relationships in one repack."""
    replacements = dict(parts)

    types = etree.fromstring(core_ooxml.read_part(source, core_ooxml.CONTENT_TYPES_PART))
    for part_name, content_type in overrides.items():
        override = etree.SubElement(types, ooxml.qn("ct:Override"))
        override.set("PartName", f"/{part_name}")
        override.set("ContentType", content_type)
    replacements[core_ooxml.CONTENT_TYPES_PART] = ooxml.serialize(types)

    for slide_number, entries in slide_rels.items():
        name = f"ppt/slides/slide{slide_number}.xml"
        rels = etree.fromstring(core_ooxml.read_part(source, ooxml.rels_path(name)))
        used = {r.get("Id") for r in rels}
        for kind, target_path in entries:
            rel_id = next(f"rId{n}" for n in range(1, 500) if f"rId{n}" not in used)
            used.add(rel_id)
            element = etree.SubElement(rels, ooxml.qn("rel:Relationship"))
            element.set("Id", rel_id)
            element.set("Type", f"{ooxml._REL_BASE}/{kind}")
            element.set("Target", target_path)
        replacements[ooxml.rels_path(name)] = ooxml.serialize(rels)

    if presentation_rels:
        rels = etree.fromstring(
            core_ooxml.read_part(source, ooxml.rels_path("ppt/presentation.xml"))
        )
        used = {r.get("Id") for r in rels}
        for kind, target_path in presentation_rels:
            rel_id = next(f"rId{n}" for n in range(1, 500) if f"rId{n}" not in used)
            used.add(rel_id)
            element = etree.SubElement(rels, ooxml.qn("rel:Relationship"))
            element.set("Id", rel_id)
            element.set("Type", f"{ooxml._REL_BASE}/{kind}")
            element.set("Target", target_path)
        replacements[ooxml.rels_path("ppt/presentation.xml")] = ooxml.serialize(rels)

    return core_ooxml.repack(source, target, replacements)


@pytest.fixture
def classic_comments_deck(simple_deck: Path, tmp_path: Path) -> Path:
    """Slides 1 and 3 carry classic comments; slide 2 has none.

    Classic comments are in scope unconditionally (spec section 7) and their
    format is stable and documented, so this fixture is trustworthy in a way the
    modern one is not.
    """
    target = tmp_path / "classic-comments.pptx"
    return _add_parts(
        simple_deck,
        target,
        {
            "ppt/commentAuthors.xml": _classic_authors(),
            "ppt/comments/comment1.xml": _classic_comments(
                [(0, "2026-01-02T03:04:05", "First thought"), (1, "2026-01-03T00:00:00", "Second")]
            ),
            "ppt/comments/comment3.xml": _classic_comments(
                [(0, "2026-02-01T09:00:00", "On the third slide")]
            ),
        },
        {
            "ppt/commentAuthors.xml": _CLASSIC_AUTHORS_TYPE,
            "ppt/comments/comment1.xml": _CLASSIC_COMMENTS_TYPE,
            "ppt/comments/comment3.xml": _CLASSIC_COMMENTS_TYPE,
        },
        {
            1: [("comments", "../comments/comment1.xml")],
            3: [("comments", "../comments/comment3.xml")],
        },
        [("commentAuthors", "commentAuthors.xml")],
    )


@pytest.fixture
def modern_comments_deck(simple_deck: Path, tmp_path: Path) -> Path:
    """A deck carrying a part with the *modern* comment content type.

    Read the module docstring before building anything on this: the content type
    is the only part of it that reflects reality. The internal markup is a
    placeholder, because spec section 11.1 forbids encoding a guess at
    PowerPoint's schema, and no reference file was available. Tests use this to
    drive the section 7 deferral, never to parse comments.
    """
    target = tmp_path / "modern-comments.pptx"
    return _add_parts(
        simple_deck,
        target,
        {"ppt/comments/modernComment_slide2.xml": b"<unverified/>"},
        {"ppt/comments/modernComment_slide2.xml": ooxml.MODERN_COMMENT_CONTENT_TYPE},
        {2: [("comments", "../comments/modernComment_slide2.xml")]},
        [],
    )


@pytest.fixture
def orphaned_modern_deck(simple_deck: Path, tmp_path: Path) -> Path:
    """A modern comment part declared in the package but attached to no slide.

    The awkward case: presence is knowable, placement is not. A reader that keys
    the deferral off "which slides have modern parts" would find none here and
    fall through to a classic-only answer — which is the silent partial result
    section 7 forbids, arrived at from the other direction.
    """
    target = tmp_path / "orphaned-modern.pptx"
    return _add_parts(
        simple_deck,
        target,
        {"ppt/comments/modernComment_orphan.xml": b"<unverified/>"},
        {"ppt/comments/modernComment_orphan.xml": ooxml.MODERN_COMMENT_CONTENT_TYPE},
        {},
        [],
    )


@pytest.fixture
def mixed_comments_deck(classic_comments_deck: Path, tmp_path: Path) -> Path:
    """Classic comments *and* a modern part — section 7's mixed case, where
    partial results are sacrificed for an error that cannot be mistaken for a
    complete read."""
    target = tmp_path / "mixed-comments.pptx"
    return _add_parts(
        classic_comments_deck,
        target,
        {"ppt/comments/modernComment_slide2.xml": b"<unverified/>"},
        {"ppt/comments/modernComment_slide2.xml": ooxml.MODERN_COMMENT_CONTENT_TYPE},
        {2: [("comments", "../comments/modernComment_slide2.xml")]},
        [],
    )


# --- odds and ends -----------------------------------------------------------


@pytest.fixture(scope="session")
def secret_text() -> str:
    return SECRET_TEXT


@pytest.fixture
def sample_image(tmp_path: Path) -> Path:
    target = tmp_path / "sample.png"
    target.write_bytes(png_bytes().getvalue())
    return target


@pytest.fixture
def emu() -> type[Emu]:
    return Emu


@pytest.fixture
def points() -> type[Pt]:
    return Pt
