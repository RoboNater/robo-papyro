"""Creating, appending, replacing, and editing decks.

Section 9's segmentation rules are the bulk of this: a document is a scroll and a
deck is a sequence, and the mapping between them has to be deterministic. Append
gets its own class because it substitutes one rule and adds two guarantees.
"""

from __future__ import annotations

import pytest

from rp_core.errors import InputError
from rp_pptx.models import CoreProperties
from rp_pptx.pptx import read, write
from rp_pptx.pptx.write import segment


class TestSegmentation:
    """Decided before any shape exists, so the rules can be read off directly."""

    def test_the_first_heading_opens_the_title_slide(self):
        plans = segment("# Deck\n", appending=False)
        assert [(p.role, p.title) for p in plans] == [("title", "Deck")]

    def test_a_paragraph_after_the_title_becomes_the_subtitle(self):
        plans = segment("# Deck\nA subtitle line\n", appending=False)
        assert plans[0].subtitle == "A subtitle line"

    def test_later_top_level_headings_open_section_slides(self):
        plans = segment("# Deck\n\n# Part Two\n\n# Part Three\n", appending=False)
        assert [p.role for p in plans] == ["title", "section", "section"]

    def test_second_level_headings_open_content_slides(self):
        plans = segment("# Deck\n\n## One\n\n## Two\n", appending=False)
        assert [(p.role, p.title) for p in plans] == [
            ("title", "Deck"),
            ("content", "One"),
            ("content", "Two"),
        ]

    def test_a_thematic_break_splits_without_changing_role(self):
        plans = segment("## One\n- a\n\n---\n\n- b\n", appending=False)
        assert [p.role for p in plans] == ["content", "content"]
        assert plans[1].title is None
        assert [b.text for b in plans[1].body] == ["b"]

    def test_a_thematic_break_between_headings_makes_no_empty_slide(self):
        plans = segment("## One\n- a\n\n---\n\n## Two\n- b\n", appending=False)
        assert len(plans) == 2

    def test_third_level_headings_are_lead_in_bullets_not_slides(self):
        """Section 9: decks do not have sub-sub-sections, outlines do."""
        plans = segment("## One\n\n### Lead in\ntext\n", appending=False)
        assert len(plans) == 1
        assert plans[0].body[0].text == "**Lead in**"

    def test_html_comments_become_notes(self):
        plans = segment("## One\n- a\n\n<!-- say this out loud -->\n", appending=False)
        assert plans[0].notes == ["say this out loud"]

    def test_list_nesting_survives(self):
        plans = segment("## One\n- top\n  - nested\n", appending=False)
        assert [(b.level, b.text) for b in plans[0].body] == [(1, "top"), (2, "nested")]

    def test_tables_and_code_are_kept_apart_from_the_body(self):
        plans = segment("## One\n\n| a |\n|---|\n| 1 |\n\n```\ncode\n```\n", appending=False)
        assert [b.kind for b in plans[0].extras] == ["table", "code"]

    def test_an_image_only_slide_takes_the_blank_role(self):
        plans = segment("![alt](pic.png)\n", appending=False)
        assert plans[0].role == "blank"

    def test_empty_markdown_makes_no_slides(self):
        assert segment("", appending=False) == []


class TestCreate:
    def test_it_builds_the_slides_the_segmentation_describes(self, tmp_path):
        write.create(
            tmp_path / "d.pptx",
            markdown="# Deck\nSub\n\n# Section\n\n## Content\n- a\n  - b\n",
        )
        index = read.get_index(tmp_path / "d.pptx")
        assert [(t.layout, t.title) for t in index.titles] == [
            ("Title Slide", "Deck"),
            ("Section Header", "Section"),
            ("Title and Content", "Content"),
        ]

    def test_bullet_levels_land_on_the_slide(self, tmp_path):
        write.create(tmp_path / "d.pptx", markdown="## C\n- a\n  - b\n    - c\n")
        paragraphs = read.get_text(tmp_path / "d.pptx")[0].paragraphs
        assert [(p.level, p.text) for p in paragraphs if p.text] == [
            (0, "C"),
            (0, "a"),
            (1, "b"),
            (2, "c"),
        ]

    def test_a_table_becomes_a_real_table(self, tmp_path):
        write.create(tmp_path / "d.pptx", markdown="## C\n\n| a | b |\n|---|---|\n| 1 | 2 |\n")
        tables = read.get_tables(tmp_path / "d.pptx")
        assert len(tables) == 1
        assert tables[0].data == [["a", "b"], ["1", "2"]]

    def test_a_comment_becomes_speaker_notes(self, tmp_path):
        write.create(tmp_path / "d.pptx", markdown="## C\n- a\n\n<!-- notes -->\n")
        assert read.get_notes(tmp_path / "d.pptx")[0].text == "notes"

    def test_an_image_is_placed(self, tmp_path, sample_image, monkeypatch):
        monkeypatch.chdir(sample_image.parent)
        write.create(tmp_path / "d.pptx", markdown=f"## C\n\n![a picture]({sample_image.name})\n")
        images = read.get_images(tmp_path / "d.pptx")
        assert len(images) == 1
        assert images[0].alt_text == "a picture"

    def test_a_missing_image_is_an_input_error_not_a_silent_skip(self, tmp_path):
        with pytest.raises(InputError, match="Image not found"):
            write.create(tmp_path / "d.pptx", markdown="## C\n\n![x](nope.png)\n")

    def test_inline_formatting_reaches_the_runs(self, tmp_path):
        write.create(tmp_path / "d.pptx", markdown="## C\n- plain **bold** and *italic*\n")
        runs = [
            run
            for slide in read.get_text(tmp_path / "d.pptx", runs=True)
            for para in slide.paragraphs
            if para.runs
            for run in para.runs
        ]
        assert any(run.bold and run.text == "bold" for run in runs)
        assert any(run.italic and run.text == "italic" for run in runs)

    def test_an_empty_deck_is_valid(self, tmp_path):
        write.create(tmp_path / "empty.pptx")
        assert read.get_index(tmp_path / "empty.pptx").slide_count == 0

    def test_the_template_deck_contributes_no_slides(self, template_env, tmp_path):
        write.create(tmp_path / "d.pptx", markdown="# T\n", template="house_like")
        assert read.get_index(tmp_path / "d.pptx").slide_count == 1


class TestPlaceholderValidation:
    """A layout can exist and still have nowhere to put the content.

    Section 5.1 makes a missing *layout* an error rather than a fallback;
    checking only the name stops one step short, because the content is then
    dropped and the deck comes out quietly missing text. Same failure, same
    treatment — and as lazy as the name check, so a slide that needs neither
    placeholder is unaffected.
    """

    @pytest.fixture
    def mapped_to_blank(self, template_env, request):
        """Point a role at ``House Blank``, which has no placeholders at all."""
        import json

        from rp_pptx import templates

        sidecar = templates.layoutmap_path(template_env / "house_like.potx")
        original = sidecar.read_text(encoding="utf-8")

        def apply(**overrides):
            mapping = json.loads(original)
            mapping.update(overrides)
            sidecar.write_text(json.dumps(mapping), encoding="utf-8")

        request.addfinalizer(lambda: sidecar.write_text(original, encoding="utf-8"))
        return apply

    def test_a_title_with_nowhere_to_go_is_an_error(self, mapped_to_blank, tmp_path):
        mapped_to_blank(content="House Blank")
        with pytest.raises(InputError) as error:
            write.create(
                tmp_path / "d.pptx", markdown="## Heading\n", template="house_like"
            )
        message = str(error.value)
        assert "House Blank" in message
        assert "title placeholder" in message

    def test_body_content_with_nowhere_to_go_is_an_error(self, mapped_to_blank, tmp_path):
        """Previously this produced a slide with the heading and every bullet
        silently discarded."""
        mapped_to_blank(content="House Blank")
        with pytest.raises(InputError) as error:
            write.create(
                tmp_path / "d.pptx",
                markdown="## Heading\n- bullet one\n- bullet two\n",
                template="house_like",
            )
        assert "House Blank" in str(error.value)

    def test_the_message_lists_what_the_layout_does_have(self, mapped_to_blank, tmp_path):
        mapped_to_blank(section="House Blank")
        with pytest.raises(InputError) as error:
            write.create(
                tmp_path / "d.pptx", markdown="# T\n\n# Section\n", template="house_like"
            )
        assert "none at all" in str(error.value)

    def test_a_layout_with_a_title_but_no_body_is_fine_until_body_arrives(
        self, mapped_to_blank, tmp_path
    ):
        """`House Section Break` has a title placeholder and no body one. A
        section slide carries no body, so it must still work."""
        write.create(tmp_path / "ok.pptx", markdown="# T\n\n# Section\n", template="house_like")
        assert [t.title for t in read.get_index(tmp_path / "ok.pptx").titles] == ["T", "Section"]

    def test_an_image_only_slide_needs_no_placeholders(
        self, template_env, tmp_path, sample_image, monkeypatch
    ):
        monkeypatch.chdir(sample_image.parent)
        write.create(
            tmp_path / "d.pptx",
            markdown=f"![a picture]({sample_image.name})\n",
            template="house_like",
        )
        assert read.get_index(tmp_path / "d.pptx").image_count == 1

    def test_the_body_placeholder_is_a_text_one_not_a_picture_one(self, tmp_path):
        """PowerPoint's "Picture with Caption" has a picture placeholder at idx 1
        and a body one at idx 2; "first placeholder with a text frame" picks the
        wrong one and puts bullets inside the picture."""
        from pptx import Presentation

        from rp_pptx.pptx.write import _body_placeholder

        presentation = Presentation()
        layout = next(
            candidate
            for candidate in presentation.slide_layouts
            if candidate.name == "Picture with Caption"
        )
        slide = presentation.slides.add_slide(layout)
        found = _body_placeholder(slide)
        assert found is not None
        assert found.placeholder_format.idx == 2


class TestAppend:
    """Section 9's three append rules, each asserted separately."""

    def test_a_leading_top_level_heading_opens_a_section_not_a_title(self, simple_deck, tmp_path):
        write.append_markdown(simple_deck, "# Added\n", output=tmp_path / "out.pptx")
        titles = read.get_index(tmp_path / "out.pptx").titles
        assert titles[-1].layout == "Section Header"
        assert titles[-1].title == "Added"

    def test_leading_unheaded_content_opens_a_new_untitled_slide(self, simple_deck, tmp_path):
        write.append_markdown(simple_deck, "just a line\n", output=tmp_path / "out.pptx")
        index = read.get_index(tmp_path / "out.pptx")
        assert index.slide_count == 4
        assert index.titles[-1].title is None

    def test_it_is_never_merged_into_the_existing_final_slide(self, simple_deck, tmp_path):
        before = read.get_text(simple_deck)[-1]
        write.append_markdown(simple_deck, "extra line\n", output=tmp_path / "out.pptx")
        after = read.get_text(tmp_path / "out.pptx")[len(read.get_text(simple_deck)) - 1]
        assert [p.text for p in after.paragraphs] == [p.text for p in before.paragraphs]

    def test_no_existing_slide_changes_at_all(self, rich_deck, tmp_path):
        before = read.get_text(rich_deck)
        notes_before = read.get_notes(rich_deck)
        write.append_markdown(rich_deck, "## New\n- x\n", output=tmp_path / "out.pptx")
        after = read.get_text(tmp_path / "out.pptx")[: len(before)]
        assert [[p.text for p in s.paragraphs] for s in after] == [
            [p.text for p in s.paragraphs] for s in before
        ]
        assert read.get_notes(tmp_path / "out.pptx")[: len(notes_before)] == notes_before

    def test_appending_requires_an_output(self, simple_deck):
        with pytest.raises(InputError):
            write.append_markdown(simple_deck, "# x\n")


class TestReplace:
    def test_it_replaces_across_slides(self, simple_deck, tmp_path):
        result = write.replace_text(simple_deck, {"Alpha": "AAA"}, output=tmp_path / "o.pptx")
        assert result.replacements == {"Alpha": 1}
        assert read.get_index(tmp_path / "o.pptx").titles[0].title == "AAA"

    def test_a_key_that_matched_nothing_reports_zero(self, simple_deck, tmp_path):
        result = write.replace_text(simple_deck, {"absent": "x"}, output=tmp_path / "o.pptx")
        assert result.replacements == {"absent": 0}

    def test_it_reaches_table_cells_and_reports_the_table(self, rich_deck, tmp_path):
        result = write.replace_text(rich_deck, {"r2c2": "CHANGED"}, output=tmp_path / "o.pptx")
        assert result.replacements == {"r2c2": 1}
        assert "table:1" in result.locations
        assert read.get_tables(tmp_path / "o.pptx")[0].data[1][1] == "CHANGED"

    def test_it_reaches_grouped_shapes(self, rich_deck, tmp_path):
        result = write.replace_text(
            rich_deck, {"text inside a group": "REPLACED"}, output=tmp_path / "o.pptx"
        )
        assert result.replacements == {"text inside a group": 1}

    def test_it_reaches_notes_slides_and_reports_them(self, rich_deck, tmp_path):
        result = write.replace_text(
            rich_deck, {"Speaker notes": "NOTES"}, output=tmp_path / "o.pptx"
        )
        assert result.replacements == {"Speaker notes": 1}
        assert "notes:1" in result.locations
        assert read.get_notes(tmp_path / "o.pptx")[0].text.startswith("NOTES")

    def test_it_does_not_touch_layouts_or_masters(self, house_like_template, tmp_path, secret_text):
        """Section 6: their text is design furniture, and editing it from a
        content operation is a surprise."""
        deck = tmp_path / "d.pptx"
        write.create(deck, markdown="# T\n", template=house_like_template)
        result = write.replace_text(deck, {secret_text: "GONE"}, output=tmp_path / "o.pptx")
        assert result.replacements == {secret_text: 0}

    def test_results_do_not_depend_on_dict_ordering(self, overlap_deck, tmp_path):
        """Section 6 rule 5, at the level a caller actually sees."""
        one = write.replace_text(
            overlap_deck, {"Hello": "1", "Hello World": "2"}, output=tmp_path / "a.pptx"
        )
        other = write.replace_text(
            overlap_deck, {"Hello World": "2", "Hello": "1"}, output=tmp_path / "b.pptx"
        )
        assert one.replacements == other.replacements
        assert [p.text for s in read.get_text(tmp_path / "a.pptx") for p in s.paragraphs] == [
            p.text for s in read.get_text(tmp_path / "b.pptx") for p in s.paragraphs
        ]

    def test_ignore_case(self, simple_deck, tmp_path):
        result = write.replace_text(
            simple_deck, {"alpha": "x"}, output=tmp_path / "o.pptx", match_case=False
        )
        assert result.replacements == {"alpha": 1}

    def test_a_replacement_spanning_runs_lands(self, runs_deck, tmp_path):
        result = write.replace_text(runs_deck, {"{{ name }}": "Ada"}, output=tmp_path / "o.pptx")
        # Title slide body, a table cell, a grouped shape, and the notes.
        assert result.replacements["{{ name }}"] == 4

    def test_replacing_requires_an_output(self, simple_deck):
        with pytest.raises(InputError):
            write.replace_text(simple_deck, {"a": "b"})


class TestSetNotes:
    def test_it_sets_the_named_slide(self, simple_deck, tmp_path):
        write.set_notes(simple_deck, 2, "note text", output=tmp_path / "o.pptx")
        assert [(n.slide_index, n.text) for n in read.get_notes(tmp_path / "o.pptx")] == [
            (2, "note text")
        ]

    @pytest.mark.parametrize("slide", [0, -1, 4, 99])
    def test_out_of_range_is_an_input_error(self, simple_deck, tmp_path, slide):
        """Index 0 used to address the *last* slide through Python's negative
        indexing — a silent wrong-target edit."""
        with pytest.raises(InputError, match="out of range"):
            write.set_notes(simple_deck, slide, "x", output=tmp_path / "o.pptx")

    def test_it_replaces_rather_than_appends(self, rich_deck, tmp_path):
        write.set_notes(rich_deck, 1, "replaced", output=tmp_path / "o.pptx")
        assert read.get_notes(tmp_path / "o.pptx")[0].text == "replaced"


class TestProperties:
    def test_setting_properties(self, simple_deck, tmp_path):
        write.set_properties(
            simple_deck, CoreProperties(title="T", author="A"), output=tmp_path / "o.pptx"
        )
        props = read.get_properties(tmp_path / "o.pptx")
        assert (props.title, props.author) == ("T", "A")
