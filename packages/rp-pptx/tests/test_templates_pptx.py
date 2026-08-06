"""Template resolution, inspection, layout maps, manifests, and synthesis.

Spec section 11.2's required assertions live here, plus the section 5 rules they
exist to protect. The three fixtures are adversarial: ``hostile`` is the one that
proves failures are loud, and that they are loud *only* when the missing role is
actually reached.
"""

from __future__ import annotations

import json

import pytest

from rp_core.errors import InputError
from rp_pptx import ooxml, templates
from rp_pptx.pptx import read, write


class TestResolution:
    def test_an_explicit_path_that_exists_is_used(self, house_like_template):
        assert templates.resolve_template(house_like_template) == house_like_template

    def test_a_bare_name_resolves_against_the_template_dir(self, template_env):
        assert templates.resolve_template("house_like").name == "house_like.potx"

    def test_potx_is_preferred_over_pptx_of_the_same_name(self, template_env, tmp_path):
        """Section 5.3: a template is what was asked for."""
        from pptx import Presentation

        both = template_env / "twins.pptx"
        Presentation().save(str(both))
        ooxml.retype_as_template(both, template_env / "twins.potx")
        assert templates.resolve_template("twins").suffix == ".potx"

    def test_a_wrong_path_reports_the_path_not_the_name(self, template_env):
        """Section 5.1 case 4. Reporting "no template called ./decks/x.potx"
        would send the user looking in the template directory, which is not
        where they meant to look at all."""
        with pytest.raises(InputError) as error:
            templates.resolve_template("./decks/missing.potx")
        assert "decks/missing.potx" in str(error.value)
        assert "available templates" not in str(error.value)

    def test_an_unresolvable_name_lists_what_is_there(self, template_env):
        with pytest.raises(InputError) as error:
            templates.resolve_template("nonexistent")
        message = str(error.value)
        assert "nonexistent" in message
        assert "house_like" in message and "hostile" in message

    def test_none_falls_back_to_the_bundled_default(self, template_env):
        assert templates.resolve_template(None) == templates.bundled_default()

    def test_the_configured_default_wins_over_the_bundled_one(self, template_env, monkeypatch):
        monkeypatch.setenv("RP_PPTX_TEMPLATE", "house_like")
        assert templates.resolve_template(None).name == "house_like.potx"

    def test_listing_finds_every_template(self, template_env):
        names = {info.name for info in templates.list_templates()}
        assert {"minimal", "house_like", "hostile"} <= names


class TestInspection:
    def test_house_like_reports_its_layouts_and_masters(self, house_like_template):
        info = templates.inspect_template(house_like_template)
        assert info.format == "potx"
        assert info.aspect_ratio == "16:9"
        assert info.master_count == 2
        names = [layout.name for layout in info.layouts]
        assert "RP Title" in names
        assert "Résumé Layout" in names, "a non-ASCII layout name must survive"

    def test_placeholder_types_are_ooxml_tokens(self, house_like_template):
        """Section 3 documents "title", "body", "pic", "tbl" — the XML tokens,
        not python-pptx's enum names, and what synthesis writes back out."""
        info = templates.inspect_template(house_like_template)
        title = next(layout for layout in info.layouts if layout.name == "RP Title")
        assert [(p.idx, p.type) for p in title.placeholders] == [
            (0, "ctrTitle"),
            (1, "subTitle"),
        ]

    def test_layout_index_is_one_based_within_its_master(self, house_like_template):
        info = templates.inspect_template(house_like_template)
        assert info.layouts[0].index == 1
        # The second master restarts at 1, which is how synthesis regroups them.
        assert info.layouts[-1].index == 1


class TestLayoutMaps:
    def test_a_sidecar_is_loaded(self, house_like_template):
        layoutmap = templates.load_layoutmap(house_like_template)
        assert layoutmap.title == "RP Title"
        assert layoutmap.content == "House Content"

    def test_no_sidecar_means_the_defaults(self, hostile_template):
        assert templates.load_layoutmap(hostile_template).title == "Title Slide"

    def test_every_default_names_a_layout_that_really_exists(self, minimal_template):
        """The rp-docx StyleMap.code lesson: a default may only name something
        that is really there. python-pptx's bundled template is the reference."""
        from rp_pptx.models import LayoutMap

        available = {layout.name for layout in templates.inspect_template(minimal_template).layouts}
        for role, name in LayoutMap().model_dump().items():
            assert name in available, f"default for {role!r} names a missing layout"

    def test_a_malformed_sidecar_is_an_input_error(self, tmp_path, minimal_template):
        import shutil

        copied = tmp_path / "copy.potx"
        shutil.copy(minimal_template, copied)
        templates.layoutmap_path(copied).write_text("{not json", encoding="utf-8")
        with pytest.raises(InputError):
            templates.load_layoutmap(copied)

    def test_the_scaffold_guesses_from_layout_names(self, house_like_template):
        guessed = templates.scaffold_layoutmap(house_like_template)
        assert guessed.title == "RP Title"
        assert guessed.section == "House Section Break"


class TestLazyLayoutChecking:
    """Section 5.1: "absent" means absent *when needed*."""

    def test_a_title_only_deck_works_on_a_template_missing_the_content_layout(
        self, template_env, tmp_path
    ):
        write.create(tmp_path / "ok.pptx", markdown="# Just a title\n", template="hostile")
        assert read.get_index(tmp_path / "ok.pptx").slide_count == 1

    def test_reaching_for_the_missing_layout_fails_loudly(self, template_env, tmp_path):
        with pytest.raises(InputError) as error:
            write.create(
                tmp_path / "bad.pptx",
                markdown="# Title\n\n## Body slide\n- point\n",
                template="hostile",
            )
        message = str(error.value)
        assert "Title and Content" in message, "names the layout it needed"
        assert "Overview" in message, "lists what the template does have"

    def test_it_never_silently_falls_back(self, template_env, tmp_path):
        """Two layouts differing only by case exist in `hostile`; neither is a
        substitute for the one that is missing."""
        with pytest.raises(InputError):
            write.create(tmp_path / "bad.pptx", markdown="## Body\n- x\n", template="hostile")


class TestAspectContract:
    """Section 4's implicit/explicit pair: the same resolved file, different
    behaviour, decided on the argument rather than the path."""

    def test_implicit_none_forces_widescreen_over_the_bundled_default(self, tmp_path):
        write.create(tmp_path / "implicit.pptx")
        assert read.get_index(tmp_path / "implicit.pptx").aspect_ratio == "16:9"

    def test_the_same_file_passed_explicitly_keeps_its_own_geometry(
        self, tmp_path, bundled_default
    ):
        write.create(tmp_path / "explicit.pptx", template=bundled_default)
        assert read.get_index(tmp_path / "explicit.pptx").aspect_ratio == "4:3"

    def test_an_explicit_template_wins_over_the_aspect_argument(self, template_env, tmp_path):
        write.create(tmp_path / "house.pptx", template="house_like", aspect="4:3")
        assert read.get_index(tmp_path / "house.pptx").aspect_ratio == "16:9"

    def test_implicit_four_three_is_honoured(self, tmp_path):
        write.create(tmp_path / "43.pptx", aspect="4:3")
        assert read.get_index(tmp_path / "43.pptx").aspect_ratio == "4:3"

    def test_an_unknown_aspect_is_rejected(self, tmp_path):
        with pytest.raises(InputError):
            write.create(tmp_path / "x.pptx", aspect="21:9")


class TestManifests:
    def test_a_manifest_carries_no_confidential_text(self, house_like_template, secret_text):
        """Section 5.2: redaction is a correctness property, not a convention.

        The template carries the secret in three places a careless manifest
        might pick up — a master text box, the author, and the title.
        """
        blob = templates.build_manifest(house_like_template).model_dump_json()
        assert secret_text not in blob

    def test_a_manifest_carries_no_path_beyond_the_basename(self, house_like_template):
        manifest = templates.build_manifest(house_like_template)
        assert manifest.name == "house_like.potx"
        assert "/" not in manifest.name

    def test_presence_flags_are_populated_not_defaulted(self, house_like_template):
        """The two fields rp-docx's equivalent checkpoint caught silently
        unpopulated (section 12 step 5)."""
        manifest = templates.build_manifest(house_like_template)
        assert manifest.master_image_count > 0, "house_like has a logo on its master"
        assert manifest.layoutmap is not None, "house_like has a layoutmap beside it"

    def test_a_manifest_round_trips_through_json(self, house_like_template, tmp_path):
        manifest = templates.build_manifest(house_like_template)
        written = templates.write_manifest(manifest, tmp_path / "m.manifest.json")
        assert json.loads(written.read_text())["name"] == "house_like.potx"
        assert templates.read_manifest(written) == manifest


class TestSynthesis:
    def test_synthesis_reproduces_the_layout_inventory(self, house_like_template, tmp_path):
        """Section 11.2's requirement, and the whole point of the manifest loop:
        CI exercises a real template's shape without the file."""
        manifest = templates.build_manifest(house_like_template)
        rebuilt = templates.synthesize(manifest, tmp_path / "rebuilt.potx")

        original = templates.inspect_template(house_like_template)
        copy = templates.inspect_template(rebuilt)

        def shape(info):
            return [
                (layout.name, layout.index, [(p.idx, p.type) for p in layout.placeholders])
                for layout in info.layouts
            ]

        assert shape(copy) == shape(original)
        assert copy.master_count == original.master_count
        assert (copy.slide_width_emu, copy.slide_height_emu) == (
            original.slide_width_emu,
            original.slide_height_emu,
        )

    def test_a_synthesized_template_is_a_template(self, house_like_template, tmp_path):
        rebuilt = templates.synthesize(
            templates.build_manifest(house_like_template), tmp_path / "r.potx"
        )
        assert ooxml.is_template(rebuilt)

    def test_a_synthesized_template_can_build_a_deck(self, house_like_template, tmp_path):
        """Structural equivalence has to mean usable, not merely readable."""
        rebuilt = templates.synthesize(
            templates.build_manifest(house_like_template), tmp_path / "r.potx"
        )
        templates.write_layoutmap(
            templates.load_layoutmap(house_like_template), templates.layoutmap_path(rebuilt)
        )
        write.create(
            tmp_path / "deck.pptx", markdown="# Title\n\n## Body\n- point\n", template=rebuilt
        )
        index = read.get_index(tmp_path / "deck.pptx")
        assert [title.layout for title in index.titles] == ["RP Title", "House Content"]

    def test_the_master_image_is_reproduced_as_a_stand_in(self, house_like_template, tmp_path):
        rebuilt = templates.synthesize(
            templates.build_manifest(house_like_template), tmp_path / "r.potx"
        )
        assert templates.build_manifest(rebuilt).master_image_count > 0


class TestContentTypes:
    """Section 5.3, asserted rather than assumed. If a future python-pptx learns
    to open templates, these fail and ``opened()`` can be simplified."""

    def test_python_pptx_still_refuses_to_open_a_potx(self, house_like_template):
        from pptx import Presentation

        with pytest.raises(ValueError, match="not a PowerPoint file"):
            Presentation(str(house_like_template))

    def test_opened_reads_a_potx_anyway(self, house_like_template):
        with ooxml.opened(house_like_template) as presentation:
            assert len(presentation.slide_masters) == 2

    def test_saving_to_potx_writes_the_template_content_type(self, tmp_path, minimal_template):
        """Without retyping this would be a mislabeled presentation, and
        PowerPoint would silently edit what the user meant to keep."""
        out = tmp_path / "saved.potx"
        write.create(out, markdown="# Title\n")
        assert ooxml.is_template(out)

    def test_saving_to_pptx_writes_the_presentation_content_type(self, tmp_path):
        out = tmp_path / "saved.pptx"
        write.create(out, markdown="# Title\n")
        assert not ooxml.is_template(out)

    def test_retyping_is_lossless_in_both_directions(self, tmp_path, minimal_template):
        there = ooxml.retype_as_presentation(minimal_template, tmp_path / "a.pptx")
        back = ooxml.retype_as_template(there, tmp_path / "b.potx")
        assert ooxml.part_names(back) == ooxml.part_names(minimal_template)
        assert ooxml.is_template(back)
