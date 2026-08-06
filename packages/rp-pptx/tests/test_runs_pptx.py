"""Run-spanning replacement (spec section 6).

Section 6 names five cases this has to cover, and they are the reason
``runs.py`` is built and unit-tested before anything depends on it: a placeholder
split across three runs, a match spanning a formatting boundary, a match inside a
table cell, a match inside a grouped shape, and overlapping candidate matches.

These run against real DrawingML built by the ``runs_deck`` fixture rather than
against stand-ins. The whole reason this module is not shared with rp-docx is
that the element shapes differ (``a:r``/``a:t`` under ``a:p``), so a test double
with a ``.runs`` list would abstract away the exact thing under test.
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from rp_pptx.ooxml import opened, qn
from rp_pptx.pptx import shapes as shape_tools
from rp_pptx.pptx.runs import find_matches, paragraph_text, replace_in_paragraph, text_spans


def split_paragraph(frame, parts: list[tuple[str, bool]]):
    """Rebuild ``frame``'s first paragraph as one run per ``(text, bold)``."""
    paragraph = frame.paragraphs[0]
    for run in list(paragraph.runs):
        run._r.getparent().remove(run._r)
    for text, bold in parts:
        run = paragraph.add_run()
        run.text = text
        run.font.bold = bold
    return paragraph


@pytest.fixture
def paragraph():
    """A standalone paragraph on a throwaway slide."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    return box.text_frame


class TestOffsetMapping:
    def test_spans_carry_their_offsets(self, paragraph):
        para = split_paragraph(paragraph, [("abc", False), ("de", False), ("f", False)])
        spans = text_spans(para)
        assert [(s.start, s.end) for s in spans] == [(0, 3), (3, 5), (5, 6)]
        assert paragraph_text(para) == "abcdef"

    def test_a_paragraph_with_no_runs_is_no_spans(self, paragraph):
        assert text_spans(paragraph.paragraphs[0]) == []
        assert replace_in_paragraph(paragraph.paragraphs[0], {"x": "y"}) == {}


class TestRequiredCases:
    """The five section 6 names, in order."""

    def test_a_placeholder_split_across_three_runs(self, paragraph):
        para = split_paragraph(
            paragraph, [("Hello {{ na", False), ("me", False), (" }} there", False)]
        )
        assert replace_in_paragraph(para, {"{{ name }}": "Ada"}) == {"{{ name }}": 1}
        assert paragraph_text(para) == "Hello Ada there"

    def test_a_match_spanning_a_formatting_boundary(self, paragraph):
        para = split_paragraph(paragraph, [("{{ na", False), ("me }}", True)])
        assert replace_in_paragraph(para, {"{{ name }}": "Ada"}) == {"{{ name }}": 1}
        assert paragraph_text(para) == "Ada"
        # Section 6 rule 4: formatting is inherited from the *first* spanned run.
        assert para.runs[0].text == "Ada"
        assert para.runs[0].font.bold is False

    def test_a_match_inside_a_table_cell(self, runs_deck):
        with opened(runs_deck) as presentation:
            table = next(
                shape.table
                for shape in shape_tools.walk(presentation.slides[1].shapes)
                if getattr(shape, "has_table", False)
            )
            cell = table.cell(0, 0)
            counts = replace_in_paragraph(cell.text_frame.paragraphs[0], {"{{ name }}": "Ada"})
            assert counts == {"{{ name }}": 1}
            assert cell.text_frame.paragraphs[0].text == "cell Ada"

    def test_a_match_inside_a_grouped_shape(self, runs_deck):
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        with opened(runs_deck) as presentation:
            # Reached through the group rather than by scanning the slide, so the
            # test cannot accidentally pass on a top-level shape.
            group = next(
                shape
                for shape in presentation.slides[1].shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP
            )
            inner = next(s for s in group.shapes if getattr(s, "has_text_frame", False))
            para = inner.text_frame.paragraphs[0]
            assert replace_in_paragraph(para, {"{{ name }}": "Ada"}) == {"{{ name }}": 1}
            assert para.text == "group Ada"

    def test_walking_reaches_grouped_shapes_at_all(self, runs_deck):
        """The traversal the case above depends on. Section 6 calls a
        replacement that misses grouped shapes the pptx version of the
        body-only bug."""
        with opened(runs_deck) as presentation:
            texts = [
                shape.text_frame.text
                for shape in shape_tools.walk(presentation.slides[1].shapes)
                if getattr(shape, "has_text_frame", False)
            ]
        assert any(text.startswith("group ") for text in texts)

    def test_overlapping_candidates_resolve_to_the_longer(self):
        """Section 6 rule 5, and the reason it exists: results must not depend on
        dict ordering, which a caller can neither see nor control."""
        one = find_matches("Hello World", {"Hello": "a", "Hello World": "b"})
        other = find_matches("Hello World", {"Hello World": "b", "Hello": "a"})
        assert [(m.key, m.start, m.end) for m in one] == [("Hello World", 0, 11)]
        assert [(m.key, m.start, m.end) for m in one] == [(m.key, m.start, m.end) for m in other]

    def test_the_shorter_key_still_matches_where_it_stands_alone(self):
        matches = find_matches("Hello World and Hello", {"Hello": "a", "Hello World": "b"})
        assert [(m.key, m.start) for m in matches] == [("Hello World", 0), ("Hello", 16)]


class TestReplacementDetails:
    def test_every_occurrence_is_replaced(self, paragraph):
        para = split_paragraph(paragraph, [("x and x and x", False)])
        assert replace_in_paragraph(para, {"x": "y"}) == {"x": 3}
        assert paragraph_text(para) == "y and y and y"

    def test_a_key_that_matches_nothing_is_absent_from_the_counts(self, paragraph):
        para = split_paragraph(paragraph, [("nothing here", False)])
        assert replace_in_paragraph(para, {"absent": "x"}) == {}

    def test_ignore_case_matches_regardless_of_case(self, paragraph):
        para = split_paragraph(paragraph, [("HELLO world", False)])
        assert replace_in_paragraph(para, {"hello": "hi"}, ignore_case=True) == {"hello": 1}
        assert paragraph_text(para) == "hi world"

    def test_case_sensitive_by_default(self, paragraph):
        para = split_paragraph(paragraph, [("HELLO world", False)])
        assert replace_in_paragraph(para, {"hello": "hi"}) == {}

    def test_a_regex_metacharacter_in_a_key_is_literal(self, paragraph):
        para = split_paragraph(paragraph, [("cost is $1.50 (net)", False)])
        assert replace_in_paragraph(para, {"$1.50 (net)": "$2.00"}) == {"$1.50 (net)": 1}
        assert paragraph_text(para) == "cost is $2.00"

    def test_surrounding_whitespace_is_preserved(self, paragraph):
        """Without ``xml:space="preserve"`` the renderer may strip the spaces and
        run the words together."""
        para = split_paragraph(paragraph, [("a KEY b", False)])
        replace_in_paragraph(para, {"KEY": " middle "})
        assert paragraph_text(para) == "a  middle  b"

    def test_preserve_formatting_off_drops_the_runs_direct_formatting(self, paragraph):
        para = split_paragraph(paragraph, [("{{ x }}", True)])
        replace_in_paragraph(para, {"{{ x }}": "plain"}, preserve_formatting=False)
        assert paragraph_text(para) == "plain"
        assert para.runs[0]._r.find(qn("a:rPr")) is None

    def test_preserve_formatting_on_keeps_it(self, paragraph):
        para = split_paragraph(paragraph, [("{{ x }}", True)])
        replace_in_paragraph(para, {"{{ x }}": "bold"}, preserve_formatting=True)
        assert para.runs[0].font.bold is True

    def test_an_empty_key_is_ignored_rather_than_matching_everywhere(self, paragraph):
        para = split_paragraph(paragraph, [("text", False)])
        assert replace_in_paragraph(para, {"": "x"}) == {}
        assert paragraph_text(para) == "text"

    def test_a_replacement_containing_the_key_does_not_loop(self, paragraph):
        para = split_paragraph(paragraph, [("a", False)])
        assert replace_in_paragraph(para, {"a": "aa"}) == {"a": 1}
        assert paragraph_text(para) == "aa"
