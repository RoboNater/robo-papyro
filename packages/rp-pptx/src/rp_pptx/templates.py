"""Template resolution, inspection, layout maps, manifests, and synthesis.

House templates are the normal path, not the exception: :func:`create` and
:func:`fill_template` default to one rather than to python-pptx's built-in
(spec section 5).

**Resolution erases explicitness, and that has a cost.** :func:`resolve_template`
returns a ``Path`` in every case, falling back to python-pptx's bundled default
rather than a ``None`` every caller would special-case. The convenience is real,
but it means a resolved path no longer says whether the caller asked for it — so
any behaviour that depends on that (today: ``create()``'s aspect rule) is decided
on the *pre-resolution* argument. Section 5.1 calls this contract, not
implementation detail.

**The manifest loop exists because real templates cannot enter this repository.**
``build_manifest`` runs against the real file wherever it lives and emits JSON
that is redacted by construction; the JSON is committed; ``synthesize``
reconstructs a structurally equivalent template at test time. CI then exercises
the real template's shape without the file ever leaving the machine holding it.
Redaction is therefore a correctness property, not a convention — a manifest
carrying slide text or an author name has failed at its one job.

Nothing here prints, and nothing here imports typer.
"""

from __future__ import annotations

import json
import os
import tempfile
from pathlib import Path
from typing import Any

import pptx
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from rp_core.errors import InputError
from rp_pptx import ooxml
from rp_pptx.models import (
    LayoutDef,
    LayoutMap,
    PlaceholderDef,
    TemplateInfo,
    TemplateManifest,
)
from rp_pptx.pptx import shapes as shape_tools

#: Suffix order matters: a `.potx` is preferred over a `.pptx` of the same name,
#: because a template is what was asked for (section 5.3).
SUFFIXES = (".potx", ".pptx")


def _roots() -> list[Path]:
    """Where a bare template name is looked up, in precedence order."""
    roots: list[Path] = []
    if configured := os.getenv("RP_PPTX_TEMPLATE_DIR"):
        roots.append(Path(configured))
    roots.extend([Path.cwd() / "templates" / "local", Path.cwd() / "templates"])
    return roots


def _looks_like_a_path(value: str) -> bool:
    """Whether the argument is a path the user got wrong, or a name to look up.

    Anything carrying a suffix or a separator was meant as a path (section 5.1
    case 4), and reporting "no template called ./decks/house.potx" for it would
    send the user looking in the wrong place entirely.
    """
    return bool(Path(value).suffix) or "/" in value or "\\" in value


def bundled_default() -> Path:
    """python-pptx's own template, the fallback when nothing is configured."""
    return Path(pptx.__file__).parent / "templates" / "default.pptx"


def available_names() -> list[str]:
    """Template names resolution would find, for an error message.

    Names only: this is the unhappy path, and opening every template found in
    order to describe one that was not found would be slow and — if any of them
    is unreadable — would replace the user's error with an unrelated one.
    """
    names: list[str] = []
    for root in _roots():
        if not root.is_dir():
            continue
        for suffix in SUFFIXES:
            names.extend(path.stem for path in sorted(root.glob(f"*{suffix}")))
    return sorted(dict.fromkeys(names))


def resolve_template(name_or_path: str | Path | None) -> Path:
    """A template path from a path, a bare name, or nothing at all."""
    if name_or_path is None:
        configured = os.getenv("RP_PPTX_TEMPLATE")
        return resolve_template(configured) if configured else bundled_default()

    candidate = Path(name_or_path)
    if candidate.is_file():
        return candidate

    raw = str(name_or_path)
    if _looks_like_a_path(raw):
        raise InputError(f"No such template file: {candidate}")

    for root in _roots():
        for suffix in SUFFIXES:
            found = root / f"{raw}{suffix}"
            if found.is_file():
                return found

    names = available_names()
    listed = ", ".join(names) if names else "none found"
    raise InputError(f"Unknown template {raw!r}; available templates: {listed}")


def list_templates() -> list[TemplateInfo]:
    """Every resolvable template, inspected. First root wins on a name clash."""
    found: dict[str, TemplateInfo] = {}
    for root in _roots():
        if not root.is_dir():
            continue
        for suffix in SUFFIXES:
            for path in sorted(root.glob(f"*{suffix}")):
                if path.stem not in found:
                    found[path.stem] = inspect_template(path)
    return list(found.values())


def _placeholder_type(placeholder: Any) -> str:
    """The OOXML token for a placeholder's type ("title", "body", "pic", …).

    The XML token rather than python-pptx's enum name, because that is what
    section 3's model documents and — more usefully — what ``synthesize`` has to
    write back out. Going through the enum's own mapping keeps the two directions
    from drifting.
    """
    kind = placeholder.placeholder_format.type
    if kind is None:
        return "body"
    try:
        return PP_PLACEHOLDER.to_xml(kind)
    except (ValueError, KeyError):
        return str(kind).split(" (")[0].lower()


def _layouts(presentation: Presentation) -> list[LayoutDef]:
    result: list[LayoutDef] = []
    for master in presentation.slide_masters:
        for index, layout in enumerate(master.slide_layouts, start=1):
            result.append(
                LayoutDef(
                    name=layout.name,
                    index=index,
                    placeholders=[
                        PlaceholderDef(
                            idx=placeholder.placeholder_format.idx,
                            type=_placeholder_type(placeholder),
                            name=placeholder.name,
                        )
                        for placeholder in layout.placeholders
                    ],
                )
            )
    return result


def inspect_template(path: Path) -> TemplateInfo:
    from rp_pptx.pptx.read import _ratio

    path = Path(path)
    with ooxml.opened(path) as presentation:
        return TemplateInfo(
            name=path.stem,
            path=path,
            format="potx" if ooxml.is_template(path) else "pptx",
            slide_width_emu=presentation.slide_width,
            slide_height_emu=presentation.slide_height,
            aspect_ratio=_ratio(presentation.slide_width, presentation.slide_height),
            master_count=len(presentation.slide_masters),
            layouts=_layouts(presentation),
        )


def layoutmap_path(template: Path) -> Path:
    """Where a template's layout map sits: ``house.potx.layoutmap.json``."""
    template = Path(template)
    return template.with_name(template.name + ".layoutmap.json")


def load_layoutmap(template: Path) -> LayoutMap:
    """The layout map beside ``template``, or the default one.

    Never hardcoded per template: house decks rarely use PowerPoint's layout
    names, and guessing which of theirs is "the content one" is exactly the
    silent-fallback section 5.1 forbids.
    """
    sidecar = layoutmap_path(template)
    if not sidecar.is_file():
        return LayoutMap()
    try:
        return LayoutMap.model_validate_json(sidecar.read_text(encoding="utf-8"))
    except ValueError as exc:
        raise InputError(f"{sidecar.name} is not a valid layout map: {exc}") from exc


def scaffold_layoutmap(template: Path) -> LayoutMap:
    """A best-effort ``LayoutMap`` guess, for a human to correct.

    Matches layout names against the words each role tends to use. A convenience
    and never authoritative — section 10 requires the CLI that emits this to say
    so. Roles it cannot guess keep the default, which will then fail loudly at
    the point of use rather than silently picking something wrong.
    """
    hints = {
        "title": ("title slide", "title", "cover", "opening"),
        "section": ("section", "divider", "break", "chapter"),
        "content": ("title and content", "content", "body", "bullet", "text"),
        "blank": ("blank", "empty", "free", "image"),
    }
    info = inspect_template(template)
    names = [layout.name for layout in info.layouts]
    guessed: dict[str, str] = {}
    for role, words in hints.items():
        for word in words:
            match = next((name for name in names if word in name.casefold()), None)
            if match:
                guessed[role] = match
                break
    return LayoutMap(**guessed)


def require_layout(presentation: Presentation, name: str) -> Any:
    """The layout called ``name``, or a loud failure naming what is there.

    **Checked at the point of use, never eagerly over the whole map** (section
    5.1). A deck with no section breaks does not need a section layout to exist,
    and an eager check would reject a perfectly usable template for a role the
    document never reaches. Never silently falls back to something else.
    """
    layouts = [layout for master in presentation.slide_masters for layout in master.slide_layouts]
    for layout in layouts:
        if layout.name == name:
            return layout
    listed = ", ".join(repr(layout.name) for layout in layouts) or "none"
    raise InputError(f"Template has no layout named {name!r}; it has: {listed}")


# --- manifests: describing a template without carrying it --------------------


def build_manifest(path: Path) -> TemplateManifest:
    """A redacted-by-construction description of a template's shape.

    Structure only — layout names, placeholder inventory, geometry, presence
    flags. Never slide text, never image bytes, never an author name, and never
    a path beyond the basename. Safe to commit, and safe to share off the
    machine holding the original.
    """
    path = Path(path)
    info = inspect_template(path)
    with ooxml.opened(path) as presentation:
        master_images = sum(
            shape_tools.is_picture(shape)
            for master in presentation.slide_masters
            for shape in master.shapes
        )
        # notes_master is lazily *created* on access, so the part list is the
        # only way to ask whether one was already there.
        notes_master = any(name.startswith("ppt/notesMasters/") for name in ooxml.part_names(path))
    sidecar = layoutmap_path(path)
    return TemplateManifest(
        name=path.name,
        format=info.format,
        slide_width_emu=info.slide_width_emu,
        slide_height_emu=info.slide_height_emu,
        aspect_ratio=info.aspect_ratio,
        master_count=info.master_count,
        layouts=info.layouts,
        master_image_count=master_images,
        notes_master_present=notes_master,
        layoutmap=load_layoutmap(path) if sidecar.is_file() else None,
    )


def _grouped_by_master(manifest: TemplateManifest) -> list[list[LayoutDef]]:
    """Layouts split back into their masters.

    ``LayoutDef.index`` is 1-based *within its master*, so a reset to 1 marks the
    start of the next master's layouts — which is how a flat list survives the
    round trip through JSON without a second key.
    """
    groups: list[list[LayoutDef]] = []
    for layout in manifest.layouts:
        if layout.index == 1 or not groups:
            groups.append([])
        groups[-1].append(layout)
    while len(groups) < manifest.master_count:
        groups.append([])
    return groups[: manifest.master_count] if manifest.master_count else groups


def synthesize(manifest: TemplateManifest, output: Path) -> Path:
    """Rebuild a structurally equivalent template from a manifest.

    Reproduces slide geometry, master count, and each layout's name and
    placeholder inventory, plus a stand-in image on the master when the original
    had one and a notes master when it was flagged. It does **not** reproduce
    theme fonts, colours, or backgrounds: this is structural equivalence for
    testing layout resolution, not visual fidelity (section 5.2).
    """
    output = Path(output)
    groups = _grouped_by_master(manifest)
    with tempfile.TemporaryDirectory(prefix="rp-pptx-synth-") as tmp:
        staged = Path(tmp) / "staged.pptx"
        presentation = Presentation()
        presentation.slide_width = manifest.slide_width_emu
        presentation.slide_height = manifest.slide_height_emu
        for _ in range(manifest.master_image_count):
            ooxml.add_master_picture(presentation, _placeholder_png(), name="Logo")
        if manifest.notes_master_present:
            # Touching it is what creates it; that is the only way in.
            presentation.notes_master  # noqa: B018
        presentation.save(str(staged))

        rebuilt = Path(tmp) / "rebuilt.pptx"
        ooxml.rebuild_masters(
            staged,
            rebuilt,
            [
                [
                    (
                        layout.name,
                        [(p.idx, p.type, p.name) for p in layout.placeholders],
                    )
                    for layout in group
                ]
                for group in groups
            ],
        )
        if output.suffix.lower() == ".potx":
            output.parent.mkdir(parents=True, exist_ok=True)
            ooxml.retype_as_template(rebuilt, output)
        else:
            output.parent.mkdir(parents=True, exist_ok=True)
            output.write_bytes(rebuilt.read_bytes())
    return output


def _placeholder_png():
    """A tiny stand-in for the logo a real master carries.

    The manifest records that an image was *there*, never the bytes — that is
    the redaction rule — so synthesis supplies its own.
    """
    from io import BytesIO

    from PIL import Image

    buffer = BytesIO()
    Image.new("RGB", (32, 32), (128, 128, 128)).save(buffer, "PNG")
    buffer.seek(0)
    return buffer


def write_manifest(manifest: TemplateManifest, output: Path) -> Path:
    """Serialize a manifest as the committed-JSON form section 5.2 describes."""
    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(manifest.model_dump_json(indent=2) + "\n", encoding="utf-8")
    return output


def read_manifest(path: Path) -> TemplateManifest:
    try:
        return TemplateManifest.model_validate_json(Path(path).read_text(encoding="utf-8"))
    except ValueError as exc:
        raise InputError(f"{Path(path).name} is not a valid template manifest: {exc}") from exc


def write_layoutmap(layoutmap: LayoutMap, output: Path) -> Path:
    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(layoutmap.model_dump(), indent=2) + "\n", encoding="utf-8")
    return output
