"""The OOXML layer: namespaces, the package zip, and the `.potx` content type.

`python-pptx` covers slides, shapes, and tables. It covers *none* of comments,
slide ordering, layout authoring, or the `.potx` content type, and those need the
package opened as what it is — a zip of XML parts. Everything reaching past
python-pptx goes through this module, so there is exactly one namespace map here
and exactly one place that knows how a `.pptx` is packed.

The mechanics underneath are :mod:`rp_core.ooxml` (spec section 12 step 2). What
stays here is everything PresentationML-specific: :data:`NS`, the content-type
strings, the part-name conventions, and the errors — core reports a missing part
as ``None`` and a malformed one as ``ValueError``, and turning those into an exit
code is a judgement only this package can make.

**The `.potx` finding, verified against python-pptx 1.0.2** (spec section 5.3,
and the `.dotx` finding from Phase 1 repeating exactly):

- ``Presentation("x.potx")`` raises ``ValueError: ... is not a PowerPoint file,
  content type is '…presentationml.template.main+xml'``. python-pptx does not
  open templates.
- ``save()`` always writes the *presentation* content type, so a file saved under
  a `.potx` name without retyping is a mislabeled presentation — PowerPoint opens
  it as an ordinary deck, silently editing what the user meant to keep as a
  template.

So retyping is load-bearing infrastructure, not a fixture convenience: every
entry point accepting a deck goes through :func:`opened`, and :func:`save`
retypes on the way out when the output is named `.potx`. ``TestContentTypes``
asserts both directions — if a future python-pptx learns to open templates, that
test fails and :func:`opened` can be simplified.

Nothing here prints, and nothing here imports typer.
"""

from __future__ import annotations

import tempfile
import zipfile
from collections.abc import Collection, Iterator
from contextlib import contextmanager
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation

from rp_core import ooxml as core_ooxml
from rp_core.errors import InputError
from rp_pptx.errors import InvalidPptxError, MissingFileError

#: Every namespace this package resolves, in one place (spec section 7).
NS: dict[str, str] = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    # Modern threaded comments and their author list (spec section 7).
    "p188": "http://schemas.microsoft.com/office/powerpoint/2018/8/main",
    "p223": "http://schemas.microsoft.com/office/powerpoint/2022/3/main",
}

PRESENTATION_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)
TEMPLATE_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
)

#: Classic per-slide comments, and the author list they index into.
CLASSIC_COMMENT_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.comments+xml"
)
CLASSIC_AUTHORS_PART = "ppt/commentAuthors.xml"

#: Modern threaded comments (spec section 7) — read as a *presence* signal only
#: while support is deferred.
MODERN_COMMENT_CONTENT_TYPE = "application/vnd.ms-powerpoint.comments+xml"
MODERN_AUTHORS_PART = "ppt/authors.xml"

PRESENTATION_PART = "ppt/presentation.xml"

#: ``xpath(element, expr)`` with :data:`NS` bound. Compiled and cached by
#: :func:`rp_core.ooxml.compiled_xpath`, which exists because python-pptx
#: overrides ``_Element.xpath`` with a version binding its own incomplete map.
xpath = core_ooxml.compiled_xpath(NS)


def qn(tag: str) -> str:
    """``"p:sldId"`` → ``"{http://…/main}sldId"`` — a Clark-notation name."""
    try:
        return core_ooxml.qualified_name(tag, NS)
    except KeyError as exc:  # a typo'd prefix is a bug here, not user input
        raise KeyError(f"{exc.args[0]}; add it to rp_pptx.ooxml.NS") from exc


def attr(element: Any, name: str, default: str | None = None) -> str | None:
    """A namespaced attribute (``"r:id"``) off ``element``."""
    return element.get(qn(name), default)


# --- the package zip ---------------------------------------------------------


def check_readable(path: Path) -> Path:
    """Fail early, and with the right error, on a path that is not a package.

    A missing file is the user's mistake (exit 1); a file that is not a zip is a
    corrupt-or-wrong-format problem (exit 3). Drawing the distinction here is
    what keeps a typo'd filename from being reported as a corrupt deck.
    """
    path = Path(path)
    if not path.exists():
        raise MissingFileError(f"No such file: {path}")
    if not path.is_file():
        raise MissingFileError(f"Not a file: {path}")
    if not zipfile.is_zipfile(path):
        raise InvalidPptxError(
            f"{path.name} is not a PowerPoint presentation: it is not an OOXML package "
            "(a .pptx/.potx file is a zip archive). A legacy binary .ppt is not supported."
        )
    return path


def part_names(path: Path) -> list[str]:
    """Every part in the package, in archive order."""
    check_readable(path)
    return core_ooxml.part_names(path)


def read_part(path: Path, name: str) -> bytes | None:
    """One part's bytes, or ``None`` when the package does not contain it."""
    check_readable(path)
    return core_ooxml.read_part(path, name)


def parse_part(path: Path, name: str) -> Any | None:
    """One part parsed as XML, or ``None`` when it is absent."""
    check_readable(path)
    try:
        return core_ooxml.parse_part(path, name)
    except ValueError as exc:
        raise InvalidPptxError(f"{path.name}: {exc}.") from exc


def repack(
    source: Path,
    target: Path,
    replacements: dict[str, bytes],
    *,
    omit: Collection[str] = (),
) -> Path:
    """Copy the package, substituting (or adding) the named parts, dropping ``omit``."""
    check_readable(source)
    return core_ooxml.repack(source, target, replacements, omit=omit)


def content_types(path: Path) -> str:
    """``[Content_Types].xml`` as text, for presence checks."""
    data = read_part(path, core_ooxml.CONTENT_TYPES_PART)
    if data is None:
        raise InvalidPptxError(
            f"{path.name} has no [Content_Types].xml; it is not an OOXML package."
        )
    return data.decode("utf-8", "replace")


# --- content types: .pptx vs .potx -------------------------------------------


def slide_parts_in_order(path: Path) -> list[str]:
    """Slide part names in **presentation order**, via ``p:sldIdLst``.

    Part *names* are not deck order and never were. Nothing renumbers
    ``slide3.xml`` when a slide moves, and this package's own ``reorder_slides``
    rewrites ``p:sldIdLst`` while leaving every part where it is — so anything
    that reads ``slide<N>.xml`` as "the Nth slide" is reading a filename that
    stopped being true the first time someone reordered a deck. Deletion breaks
    it a second way, by leaving the surviving numbers non-contiguous.

    The relationship graph is the only source of truth, so this walks it.
    """
    presentation = parse_part(path, PRESENTATION_PART)
    rels = parse_part(path, rels_path(PRESENTATION_PART))
    if presentation is None or rels is None:
        return []
    targets = {entry.get("Id"): entry.get("Target") for entry in rels}
    ordered: list[str] = []
    for entry in xpath(presentation, "./p:sldIdLst/p:sldId"):
        target = targets.get(attr(entry, "r:id"))
        if target:
            ordered.append(core_ooxml.resolve_target(PRESENTATION_PART, target))
    return ordered


def related_parts(path: Path, part: str, kind: str) -> list[str]:
    """Parts ``part`` points at through a relationship ending in ``kind``."""
    rels = parse_part(path, rels_path(part))
    if rels is None:
        return []
    return [
        core_ooxml.resolve_target(part, relationship.get("Target"))
        for relationship in rels
        if relationship.get("Type", "").endswith(f"/{kind}") and relationship.get("Target")
    ]


def comment_parts_by_slide(path: Path) -> dict[int, list[tuple[str, str]]]:
    """1-based presentation slide number → its ``(part, content_type)`` comments.

    Classic and modern comment parts share a directory and a naming convention,
    so the content type is what tells them apart — not the filename, and not the
    relationship type, which is the piece of the modern format that could not be
    verified against a real PowerPoint file (spec section 7).
    """
    declared = core_ooxml.override_content_types(path)
    found: dict[int, list[tuple[str, str]]] = {}
    for number, slide_part in enumerate(slide_parts_in_order(path), start=1):
        attached = [
            (part, declared.get(part, ""))
            for part in related_parts(path, slide_part, "comments")
        ]
        if attached:
            found[number] = attached
    return found


def is_template(path: Path) -> bool:
    """Whether the package declares itself a template (`.potx`)."""
    return TEMPLATE_CONTENT_TYPE in content_types(path)


def has_modern_comments(path: Path) -> bool:
    """Whether the package carries modern threaded-comment parts (section 7)."""
    return MODERN_COMMENT_CONTENT_TYPE in content_types(path)


def retype_as_presentation(path: Path, output: Path | None = None) -> Path:
    """Rewrite the main content type to the presentation one (`.pptx`)."""
    check_readable(path)
    return core_ooxml.retype(path, output, TEMPLATE_CONTENT_TYPE, PRESENTATION_CONTENT_TYPE)


def retype_as_template(path: Path, output: Path | None = None) -> Path:
    """Rewrite the main content type to the template one (`.potx`)."""
    check_readable(path)
    return core_ooxml.retype(path, output, PRESENTATION_CONTENT_TYPE, TEMPLATE_CONTENT_TYPE)


# --- authoring masters and layouts -------------------------------------------
#
# python-pptx can *read* layouts and *rename* them, but it cannot create one, and
# it cannot add a master at all. Spec section 5.2 needs both: synthesize() has to
# rebuild a house template's layout inventory from a manifest. So this is the
# raw-XML corner section 7 anticipates, kept here with the namespace map rather
# than scattered through templates.py.

_XML_HEADER = b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'

SLIDE_MASTER_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"
)
SLIDE_LAYOUT_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"
)
_REL_BASE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def serialize(element: Any) -> bytes:
    """An XML part's bytes, with the declaration OPC readers expect."""
    return _XML_HEADER + etree.tostring(element, encoding="UTF-8")


def rels_path(part_name: str) -> str:
    """``"ppt/presentation.xml"`` → ``"ppt/_rels/presentation.xml.rels"``."""
    head, _, tail = part_name.rpartition("/")
    return f"{head}/_rels/{tail}.rels" if head else f"_rels/{tail}.rels"


def build_rels(entries: list[tuple[str, str, str]]) -> bytes:
    """A ``.rels`` part from ``(id, type_suffix, target)`` triples."""
    root = etree.Element(qn("rel:Relationships"), nsmap={None: NS["rel"]})
    for rel_id, kind, target in entries:
        element = etree.SubElement(root, qn("rel:Relationship"))
        element.set("Id", rel_id)
        element.set("Type", f"{_REL_BASE}/{kind}")
        element.set("Target", target)
    return serialize(root)


def placeholder_shape(idx: int, ph_type: str, name: str, shape_id: int) -> Any:
    """One empty placeholder ``p:sp``, as a layout declares it.

    Only what layout inheritance keys on: the ``idx``, the ``type``, and the
    shape name. No geometry — a placeholder without an ``a:xfrm`` inherits its
    position from the master, which is exactly what a real layout does for the
    placeholders it does not move.
    """
    shape = etree.Element(qn("p:sp"))
    nv = etree.SubElement(shape, qn("p:nvSpPr"))
    cnv = etree.SubElement(nv, qn("p:cNvPr"))
    cnv.set("id", str(shape_id))
    cnv.set("name", name)
    sp_pr = etree.SubElement(nv, qn("p:cNvSpPr"))
    etree.SubElement(sp_pr, qn("a:spLocks")).set("noGrp", "1")
    nv_pr = etree.SubElement(nv, qn("p:nvPr"))
    placeholder = etree.SubElement(nv_pr, qn("p:ph"))
    # A `type` of "body" is the schema default and PowerPoint omits it; writing
    # it explicitly is still valid and keeps the round trip readable.
    placeholder.set("type", ph_type)
    placeholder.set("idx", str(idx))
    etree.SubElement(shape, qn("p:spPr"))
    text = etree.SubElement(shape, qn("p:txBody"))
    etree.SubElement(text, qn("a:bodyPr"))
    etree.SubElement(text, qn("a:lstStyle"))
    etree.SubElement(text, qn("a:p"))
    return shape


def retitle_layout(layout_xml: bytes, name: str, placeholders: list[tuple[int, str, str]]) -> bytes:
    """A copy of ``layout_xml`` renamed, carrying exactly ``placeholders``.

    Built by rewriting a real layout part rather than authoring one from
    scratch, so the colour map override, the theme link, and the schema's
    element ordering all come along — the parts that are tedious to get right
    and invisible when wrong.
    """
    layout = etree.fromstring(layout_xml)
    common = layout.find(qn("p:cSld"))
    common.set("name", name)
    tree = common.find(qn("p:spTree"))
    for child in tree:
        # Keep the group's own properties; drop every shape the source had.
        if child.tag not in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            tree.remove(child)
    for offset, (idx, ph_type, ph_name) in enumerate(placeholders):
        tree.append(placeholder_shape(idx, ph_type, ph_name, shape_id=offset + 2))
    return serialize(layout)


def add_master_picture(presentation: Presentation, image: Any, *, name: str = "Logo") -> str:
    """Put a picture on the first slide master, returning its relationship id.

    ``master.shapes`` has no ``add_picture`` — that lives on slide shape trees
    only — so the image part is registered through the master part's own API and
    the ``p:pic`` is appended by hand.
    """
    master = presentation.slide_masters[0]
    image_part, rel_id = master.part.get_or_add_image_part(image)
    pic = etree.SubElement(master.shapes._spTree, qn("p:pic"))
    nv = etree.SubElement(pic, qn("p:nvPicPr"))
    cnv = etree.SubElement(nv, qn("p:cNvPr"))
    cnv.set("id", "1023")
    cnv.set("name", name)
    cnv.set("descr", name)
    etree.SubElement(nv, qn("p:cNvPicPr"))
    etree.SubElement(nv, qn("p:nvPr"))
    fill = etree.SubElement(pic, qn("p:blipFill"))
    etree.SubElement(fill, qn("a:blip")).set(qn("r:embed"), rel_id)
    etree.SubElement(etree.SubElement(fill, qn("a:stretch")), qn("a:fillRect"))
    shape_pr = etree.SubElement(pic, qn("p:spPr"))
    xfrm = etree.SubElement(shape_pr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", "228600")
    off.set("y", "228600")
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", "685800")
    ext.set("cy", "514350")
    geom = etree.SubElement(shape_pr, qn("a:prstGeom"))
    geom.set("prst", "rect")
    etree.SubElement(geom, qn("a:avLst"))
    del image_part  # registered by get_or_add_image_part; nothing else to do
    return rel_id


#: One layout as ``rebuild_masters`` takes it: a name plus ``(idx, type, name)``
#: placeholder triples, where ``type`` is the OOXML token ("title", "body",
#: "pic", "tbl", …) that :meth:`PP_PLACEHOLDER.to_xml` produces.
LayoutSpec = tuple[str, list[tuple[int, str, str]]]


def rebuild_masters(source: Path, target: Path, masters: list[list[LayoutSpec]]) -> Path:
    """Rewrite ``source`` so it has exactly ``masters`` masters and their layouts.

    Every master is a copy of the source's first master — its colour map, theme
    link, and text styles come along — carrying a freshly authored layout list.

    **Only valid on a package with no slides.** A slide holds a relationship to
    the layout it was built from, and this replaces every layout part, so a deck
    with slides would come out pointing at parts that no longer exist. Both
    callers (``synthesize`` and the test fixtures) build templates, which have
    none, and that is the constraint that keeps this from needing to rewrite
    slide rels too.
    """
    source, target = Path(source), Path(target)
    master_xml = read_part(source, "ppt/slideMasters/slideMaster1.xml")
    layout_xml = read_part(source, "ppt/slideLayouts/slideLayout1.xml")
    if master_xml is None or layout_xml is None:
        raise InvalidPptxError(f"{source.name} has no slide master to build from.")

    theme_target = _first_rel_target(source, "ppt/slideMasters/slideMaster1.xml", "theme")
    replacements: dict[str, bytes] = {}
    layout_parts: list[str] = []
    master_parts: list[str] = []
    layout_number = 0

    for master_index, layouts in enumerate(masters, start=1):
        master_part = f"ppt/slideMasters/slideMaster{master_index}.xml"
        master_parts.append(master_part)
        rels: list[tuple[str, str, str]] = []
        master = etree.fromstring(master_xml)
        id_list = master.find(qn("p:sldLayoutIdLst"))
        for child in list(id_list):
            id_list.remove(child)

        for offset, (name, placeholders) in enumerate(layouts, start=1):
            layout_number += 1
            layout_part = f"ppt/slideLayouts/slideLayout{layout_number}.xml"
            layout_parts.append(layout_part)
            rel_id = f"rId{offset}"
            entry = etree.SubElement(id_list, qn("p:sldLayoutId"))
            # Layout ids must be unique across the package and above 2^31.
            entry.set("id", str(2147483648 + layout_number))
            entry.set(qn("r:id"), rel_id)
            rels.append((rel_id, "slideLayout", f"../slideLayouts/slideLayout{layout_number}.xml"))
            replacements[layout_part] = retitle_layout(layout_xml, name, placeholders)
            replacements[rels_path(layout_part)] = build_rels(
                [("rId1", "slideMaster", f"../slideMasters/slideMaster{master_index}.xml")]
            )

        rels.append((f"rId{len(layouts) + 1}", "theme", theme_target))
        replacements[master_part] = serialize(master)
        replacements[rels_path(master_part)] = build_rels(rels)

    _rewrite_presentation_masters(source, master_parts, replacements)
    _rewrite_content_types(source, master_parts, layout_parts, replacements)

    stale = {
        name
        for name in part_names(source)
        if name.startswith(("ppt/slideLayouts/", "ppt/slideMasters/"))
    }
    return repack(source, target, replacements, omit=stale - set(replacements))


def _first_rel_target(path: Path, part: str, kind: str) -> str:
    rels = parse_part(path, rels_path(part))
    for relationship in rels if rels is not None else []:
        if relationship.get("Type", "").endswith(f"/{kind}"):
            return relationship.get("Target")
    raise InvalidPptxError(f"{path.name}: {part} has no {kind} relationship.")


def _rewrite_presentation_masters(
    source: Path, master_parts: list[str], replacements: dict[str, bytes]
) -> None:
    """Point ``p:sldMasterIdLst`` at the new masters, keeping every other rel."""
    rels = parse_part(source, rels_path(PRESENTATION_PART))
    kept = [r for r in rels if not r.get("Type", "").endswith("/slideMaster")]
    used = {r.get("Id") for r in kept}
    entries = [(r.get("Id"), r.get("Type").rsplit("/", 1)[-1], r.get("Target")) for r in kept]

    assigned: list[str] = []
    for part in master_parts:
        rel_id = next(f"rId{n}" for n in range(1, 500) if f"rId{n}" not in used)
        used.add(rel_id)
        assigned.append(rel_id)
        entries.append((rel_id, "slideMaster", part.removeprefix("ppt/")))
    replacements[rels_path(PRESENTATION_PART)] = build_rels(entries)

    presentation = etree.fromstring(read_part(source, PRESENTATION_PART))
    id_list = presentation.find(qn("p:sldMasterIdLst"))
    for child in list(id_list):
        id_list.remove(child)
    for offset, rel_id in enumerate(assigned, start=1):
        entry = etree.SubElement(id_list, qn("p:sldMasterId"))
        entry.set("id", str(2147483648 + offset))
        entry.set(qn("r:id"), rel_id)
    replacements[PRESENTATION_PART] = serialize(presentation)


def _rewrite_content_types(
    source: Path, master_parts: list[str], layout_parts: list[str], replacements: dict[str, bytes]
) -> None:
    """Declare exactly the master and layout parts that now exist."""
    types = etree.fromstring(read_part(source, core_ooxml.CONTENT_TYPES_PART))
    for override in list(types):
        name = override.get("PartName") or ""
        if name.startswith(("/ppt/slideLayouts/", "/ppt/slideMasters/")):
            types.remove(override)
    for part, content_type in [(p, SLIDE_MASTER_CONTENT_TYPE) for p in master_parts] + [
        (p, SLIDE_LAYOUT_CONTENT_TYPE) for p in layout_parts
    ]:
        override = etree.SubElement(types, qn("ct:Override"))
        override.set("PartName", f"/{part}")
        override.set("ContentType", content_type)
    replacements[core_ooxml.CONTENT_TYPES_PART] = serialize(types)


@contextmanager
def opened(path: Path) -> Iterator[Presentation]:
    """A python-pptx ``Presentation`` for a `.pptx` **or** a `.potx`.

    A template is retyped into a temp copy first, because python-pptx refuses to
    open one at all (see the module docstring).

    **The ``yield`` is deliberately outside the ``try``.** Guarding it would make
    this context manager catch every exception raised by the caller's ``with``
    body and re-raise it as a corrupt-file error — so a bad ``--slides`` spec
    would report exit 3 on a perfectly healthy deck, and a bug in read code would
    arrive disguised as a damaged file. Only the open itself is ours to explain.
    """
    path = check_readable(path)
    with tempfile.TemporaryDirectory(prefix="rp-pptx-open-") as tmp:
        target = path
        if is_template(path):
            target = retype_as_presentation(path, Path(tmp) / "template.pptx")
        try:
            presentation = Presentation(str(target))
        except Exception as exc:
            raise InvalidPptxError(f"Cannot open {path.name}: {exc}") from exc
        yield presentation


def save(presentation: Presentation, output: Path) -> Path:
    """Save, retyping to a template when ``output`` is named `.potx`.

    python-pptx always writes the presentation content type, so without this a
    `.potx` would be a mislabeled deck — the failure section 5.3 exists to
    prevent.
    """
    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    if output.suffix.lower() == ".potx":
        with tempfile.TemporaryDirectory(prefix="rp-pptx-save-") as tmp:
            staged = Path(tmp) / "deck.pptx"
            presentation.save(str(staged))
            retype_as_template(staged, output)
    else:
        presentation.save(str(output))
    return output


def copy_for_edit(path: Path, output: Path | None) -> tuple[Presentation, Path]:
    """An open presentation detached from its source, plus where it will land.

    Detached because the source may be a temp retype of a `.potx` that
    :func:`opened` is about to delete, and because an edit must never write
    through to its input by accident. ``output`` is required: refusing to guess
    a filename is spec section 10's rule, and the CLI is what turns
    ``--in-place`` into an explicit path.

    Reading the staged copy back out of a temp directory is safe because
    python-pptx deserializes the whole package eagerly — it keeps no handle on
    the file it was opened from.
    """
    if output is None:
        raise InputError(
            "An output path is required — this package never overwrites implicitly. "
            "Pass output=... (the CLI spells this -o OUT or --in-place)."
        )
    with tempfile.TemporaryDirectory(prefix="rp-pptx-edit-") as tmp:
        staged = Path(tmp) / "staged.pptx"
        with opened(path) as source:
            source.save(str(staged))
        return Presentation(str(staged)), Path(output)
